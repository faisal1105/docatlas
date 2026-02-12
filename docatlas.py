#!/usr/bin/env python3
"""
DocAtlas document processing pipeline:
- Extract text from PDF/DOCX/PPTX/XLSX
- Summarize, categorize, tag with Azure OpenAI
- Detect duplicates via hashes + embeddings
- Output Excel files (peers + full_text)
- Move files into category folders (duplicates to <category>_Duplicate)
"""
from __future__ import annotations

import argparse
import io
import hashlib
import json
import logging
import os
import re
import shutil
import sys
import time
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple
from concurrent.futures import ThreadPoolExecutor, as_completed
import threading
import queue

import numpy as np
import pandas as pd
import requests

try:
    import pdfplumber
except Exception:  # pragma: no cover
    pdfplumber = None

try:
    import docx
except Exception:  # pragma: no cover
    docx = None

try:
    import pptx
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:  # pragma: no cover
    pptx = None
    MSO_SHAPE_TYPE = None

try:
    import openpyxl
except Exception:  # pragma: no cover
    openpyxl = None

try:
    import pytesseract
except Exception:  # pragma: no cover
    pytesseract = None

try:
    from PIL import Image
except Exception:  # pragma: no cover
    Image = None

try:
    from pdf2image import convert_from_path
except Exception:  # pragma: no cover
    convert_from_path = None

try:
    import ocrmypdf
except Exception:  # pragma: no cover
    ocrmypdf = None

try:
    import tkinter as tk
    from tkinter import filedialog, messagebox
    from tkinter import ttk
except Exception:  # pragma: no cover
    tk = None
    messagebox = None
    ttk = None

try:
    from tqdm import tqdm
except Exception:  # pragma: no cover
    tqdm = None


SUPPORTED_EXTS = {".pdf", ".doc", ".docx", ".pptx", ".xlsx"}
INVALID_WIN_CHARS = r'<>:"/\\|?*'
ILLEGAL_EXCEL_CHARS_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")

DEFAULT_CHAT_BASE_URL = "https://api.geneai.thermofisher.com/dev/gpt5"
DEFAULT_EMBEDDINGS_BASE_URL = "https://api.geneai.thermofisher.com/dev/embeddingsv2"
DEFAULT_API_VERSION = "2025-03-01-preview"
DEFAULT_CHAT_DEPLOYMENT = "gpt-5.2"
DEFAULT_EMBEDDINGS_DEPLOYMENT = "text-embedding-3-small"
DEFAULT_API_KEY_HEADER = "api-key"
DEFAULT_CHAT_PATH = "/openai/deployments/{deployment}/chat/completions"
DEFAULT_EMBEDDINGS_PATH = "/openai/deployments/{deployment}/embeddings"
DEFAULT_API_DELAY_SEC = 0.3

MAX_CHARS_PER_CHUNK = 12000
MAX_ARTICLE_CHARS = 20000
DUPLICATE_THRESHOLD = 0.97
MIN_EXTRACTED_CHARS = 200
MIN_ARTICLE_BODY_CHARS = 400
MIN_HEADING_GAP_LINES = 3
MIN_EMBEDDING_CHARS = 500
MIN_EMBEDDING_CHARS_SUMMARY = 200
MAX_TAGS = 10
RESUME_FILENAME = "resume.json"
LAST_RUN_STATS_FILENAME = "last_run_stats.json"
DEFAULT_EMBEDDINGS_SOURCE = "full_text"
DEFAULT_ESTIMATE_SEC_PER_FILE = 50.0
DEFAULT_ESTIMATE_SEC_PER_MB = 1.5
EMBEDDINGS_SOURCE_NONE = "none"
UNREADABLE_CATEGORY = "Unreadable"

USAGE_LOCK = threading.Lock()
USAGE: Dict[str, int] = {"chat_in": 0, "chat_out": 0, "embed_in": 0}

THEME = {
    # One Half Light-inspired
    "bg": "#FAFAFA",
    "panel": "#FFFFFF",
    "text_bg": "#FFFFFF",
    "fg": "#383A42",
    "muted": "#A0A1A7",
    "accent": "#61AFEF",
    "accent_dark": "#4A9FE0",
    "border": "#E5E5E5",
    "btn_bg": "#EFEFEF",
    "btn_fg": "#383A42",
}

FONT_FAMILY = "Cascadia Mono"
FONT_BASE_SIZE = 11
FONT_BASE = (FONT_FAMILY, FONT_BASE_SIZE)
FONT_SMALL = (FONT_FAMILY, FONT_BASE_SIZE - 1)
FONT_LABEL = (FONT_FAMILY, FONT_BASE_SIZE)
FONT_HEADER = (FONT_FAMILY, FONT_BASE_SIZE + 2, "bold")
FONT_BUTTON = (FONT_FAMILY, FONT_BASE_SIZE)


@dataclass
class AzureConfig:
    api_key: str
    chat_api_key: str
    embeddings_api_key: str
    api_version: str
    api_key_header: str
    chat_base_url: str
    chat_path: str
    chat_deployment: str
    embeddings_base_url: str
    embeddings_path: str
    embeddings_deployment: str
    include_model_in_body: bool


@dataclass
class DocRecord:
    doc_id: str
    file_key: str
    file_name: str
    file_path: str
    file_ext: str
    category: str
    tags: List[str]
    short_summary: str
    long_summary: str
    word_count: int
    char_count: int
    extraction_status: str
    review_flags: str
    duplicate_of: str
    duplicate_score: Optional[float]
    duplicate_group_id: str
    moved_to: str


@dataclass
class ArticleRecord:
    doc_id: str
    file_key: str
    file_name: str
    file_path: str
    article_index: int
    article_title: str
    article_summary: str
    duplicate_of: str
    duplicate_score: Optional[float]
    duplicate_group_id: str


def env(name: str, default: Optional[str] = None) -> Optional[str]:
    val = os.getenv(name)
    return val if val not in (None, "") else default


def reset_usage() -> None:
    with USAGE_LOCK:
        USAGE["chat_in"] = 0
        USAGE["chat_out"] = 0
        USAGE["embed_in"] = 0


def add_chat_usage(in_chars: int, out_chars: int) -> None:
    with USAGE_LOCK:
        USAGE["chat_in"] += in_chars
        USAGE["chat_out"] += out_chars


def add_embed_usage(in_chars: int) -> None:
    with USAGE_LOCK:
        USAGE["embed_in"] += in_chars


def get_usage() -> Dict[str, int]:
    with USAGE_LOCK:
        return dict(USAGE)


def resolve_embeddings_source(value: Optional[str]) -> str:
    if value in ("summary", "full_text", EMBEDDINGS_SOURCE_NONE):
        return value
    return DEFAULT_EMBEDDINGS_SOURCE


def apply_theme(root: tk.Tk) -> None:
    root.configure(bg=THEME["bg"])
    try:
        root.option_add("*Font", FONT_BASE)
    except Exception:
        pass
    if ttk is not None:
        try:
            style = ttk.Style(root)
            style.theme_use("clam")
            style.configure(
                "TProgressbar",
                troughcolor=THEME["panel"],
                background=THEME["accent"],
                bordercolor=THEME["border"],
                lightcolor=THEME["accent"],
                darkcolor=THEME["accent"],
            )
        except Exception:
            pass


def load_app_config(path: Path) -> Dict[str, List[str]]:
    if not path.exists():
        return {}
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {}
    apps = data.get("applications", {})
    if isinstance(apps, dict):
        cleaned: Dict[str, List[str]] = {}
        for k, v in apps.items():
            if isinstance(v, list):
                cleaned[k] = [str(x).strip() for x in v if str(x).strip()]
        return cleaned
    return {}


def save_app_config(path: Path, apps: Dict[str, List[str]]) -> None:
    data = {"applications": apps}
    path.write_text(json.dumps(data, indent=2), encoding="utf-8")


def check_ocr_dependencies(ocrmypdf_enabled: bool) -> List[str]:
    missing: List[str] = []
    if ocrmypdf_enabled:
        if shutil.which("tesseract") is None:
            missing.append("tesseract")
        if shutil.which("qpdf") is None:
            missing.append("qpdf")
        # Ghostscript executable name varies on Windows
        if shutil.which("gswin64c") is None and shutil.which("gswin32c") is None and shutil.which("gs") is None:
            missing.append("ghostscript")
    # Poppler tools for pdf2image fallback
    if shutil.which("pdftoppm") is None:
        missing.append("poppler(pdftoppm)")
    return missing


def warn_missing_ocr_deps(ocrmypdf_enabled: bool) -> List[str]:
    missing = check_ocr_dependencies(ocrmypdf_enabled)
    if not missing:
        return []
    msg = "Missing OCR dependencies (OCR may be limited): " + ", ".join(missing)
    logging.warning(msg)
    if tk is not None and messagebox is not None:
        messagebox.showwarning("OCR Dependencies", msg)
    return missing


def sanitize_folder(name: str) -> str:
    name = (name or "").strip()
    if not name:
        return "uncategorized"
    name = "".join("_" if c in INVALID_WIN_CHARS else c for c in name)
    name = name.strip(" .")
    return name or "uncategorized"


def setup_logging(out_dir: Path) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    log_path = out_dir / "docatlas.log"
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s [%(levelname)s] %(message)s",
        handlers=[logging.FileHandler(log_path, encoding="utf-8"), logging.StreamHandler(sys.stdout)],
    )


def build_url(base: str, path_template: str, deployment: str) -> str:
    if path_template.startswith("http://") or path_template.startswith("https://"):
        return path_template.format(deployment=deployment)
    base = base.rstrip("/")
    path = path_template.format(deployment=deployment)
    return f"{base}{path}"


def azure_config_from_env(require_key: bool = True) -> AzureConfig:
    api_key = env("AZURE_OPENAI_API_KEY", "")
    if require_key and not api_key and not env("AZURE_CHAT_API_KEY") and not env("AZURE_EMBEDDINGS_API_KEY"):
        raise ValueError("AZURE_OPENAI_API_KEY is not set")

    return AzureConfig(
        api_key=api_key,
        chat_api_key=env("AZURE_CHAT_API_KEY", api_key),
        embeddings_api_key=env("AZURE_EMBEDDINGS_API_KEY", api_key),
        api_version=env("AZURE_OPENAI_API_VERSION", DEFAULT_API_VERSION),
        api_key_header=env("AZURE_OPENAI_API_KEY_HEADER", DEFAULT_API_KEY_HEADER),
        chat_base_url=env("AZURE_CHAT_BASE_URL", DEFAULT_CHAT_BASE_URL),
        chat_path=env("AZURE_CHAT_PATH", DEFAULT_CHAT_PATH),
        chat_deployment=env("AZURE_CHAT_DEPLOYMENT", DEFAULT_CHAT_DEPLOYMENT),
        embeddings_base_url=env("AZURE_EMBEDDINGS_BASE_URL", DEFAULT_EMBEDDINGS_BASE_URL),
        embeddings_path=env("AZURE_EMBEDDINGS_PATH", DEFAULT_EMBEDDINGS_PATH),
        embeddings_deployment=env("AZURE_EMBEDDINGS_DEPLOYMENT", DEFAULT_EMBEDDINGS_DEPLOYMENT),
        include_model_in_body=(env("AZURE_INCLUDE_MODEL_IN_BODY", "1") == "1"),
    )


def api_delay_sec() -> float:
    val = env("DOCATLAS_API_DELAY", "")
    try:
        if val:
            return max(0.0, float(val))
    except Exception:
        pass
    return DEFAULT_API_DELAY_SEC


def list_files(input_dir: Path) -> List[Path]:
    files: List[Path] = []
    for p in input_dir.rglob("*"):
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXTS:
            files.append(p)
    return files


def scan_input_stats(files: List[Path]) -> Dict[str, Any]:
    total_size = 0
    by_ext: Dict[str, int] = {}
    for p in files:
        try:
            total_size += p.stat().st_size
        except Exception:
            pass
        ext = p.suffix.lower()
        by_ext[ext] = by_ext.get(ext, 0) + 1
    total_size_mb = total_size / (1024 * 1024)
    return {"count": len(files), "total_size_mb": total_size_mb, "by_ext": by_ext}


def load_last_run_stats(out_dir: Path) -> Optional[Dict[str, Any]]:
    stats_path = out_dir / LAST_RUN_STATS_FILENAME
    if not stats_path.exists():
        return None
    try:
        with stats_path.open("r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return None


def save_last_run_stats(out_dir: Path, stats: Dict[str, Any]) -> None:
    stats_path = out_dir / LAST_RUN_STATS_FILENAME
    try:
        with stats_path.open("w", encoding="utf-8") as f:
            json.dump(stats, f, indent=2)
    except Exception:
        pass


def format_duration(seconds: float) -> str:
    seconds = max(0, int(seconds))
    minutes = seconds / 60
    hours = minutes / 60
    days = hours / 24
    if days >= 1:
        return f"{days:.1f} days"
    if hours >= 1:
        return f"{hours:.1f} hours"
    if minutes >= 1:
        return f"{minutes:.1f} minutes"
    return f"{seconds} seconds"


def quick_estimate_runtime(
    input_stats: Dict[str, Any],
    output_dir: Path,
    ocrmypdf_enabled: bool,
    embeddings_source: str,
    chat_deployment: str,
) -> Tuple[Optional[float], str, bool]:
    baseline = load_last_run_stats(output_dir)
    count = input_stats.get("count", 0)
    total_mb = input_stats.get("total_size_mb", 0.0)
    if count <= 0:
        return None, "none", False

    if baseline and baseline.get("processed_files", 0) > 0:
        sec_per_file = baseline["elapsed_sec"] / baseline["processed_files"]
        sec_per_mb = 0.0
        if baseline.get("total_size_mb", 0) > 0:
            sec_per_mb = baseline["elapsed_sec"] / baseline["total_size_mb"]
        est_by_file = sec_per_file * count
        est_by_mb = sec_per_mb * total_mb if sec_per_mb > 0 else 0.0
        est_sec = max(est_by_file, est_by_mb)
        settings_match = (
            baseline.get("ocr_enabled") == ocrmypdf_enabled
            and baseline.get("embeddings_source") == embeddings_source
            and baseline.get("chat_deployment") == chat_deployment
        )
        logging.info(
            "Quick estimate (baseline): files=%d, size=%.1f MB -> ~%ds",
            count,
            total_mb,
            int(est_sec),
        )
        if not settings_match:
            logging.info("Quick estimate note: baseline settings differ from current run.")
        return est_sec, "baseline", settings_match

    est_by_file = DEFAULT_ESTIMATE_SEC_PER_FILE * count
    est_by_mb = DEFAULT_ESTIMATE_SEC_PER_MB * total_mb
    est_sec = max(est_by_file, est_by_mb)
    logging.info(
        "Quick estimate (heuristic): files=%d, size=%.1f MB -> ~%ds",
        count,
        total_mb,
        int(est_sec),
    )
    return est_sec, "heuristic", True


def normalize_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text or "").strip().lower()
    return text


def hash_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8", errors="ignore")).hexdigest()


def embedding_text_for_doc(embeddings_source: str, normalized_text: str, long_summary: str, short_summary: str) -> str:
    if embeddings_source == "full_text":
        return normalized_text
    summary = long_summary or short_summary or ""
    return normalize_text(summary)


def embedding_text_for_article(embeddings_source: str, body: str, article_summary: str) -> str:
    if embeddings_source == "full_text":
        return normalize_text(body)
    summary = article_summary or ""
    return normalize_text(summary)


def min_embedding_chars_for_source(embeddings_source: str) -> int:
    return MIN_EMBEDDING_CHARS if embeddings_source == "full_text" else MIN_EMBEDDING_CHARS_SUMMARY


def split_text(text: str, max_chars: int) -> List[str]:
    text = text or ""
    if len(text) <= max_chars:
        return [text]
    chunks: List[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + max_chars)
        chunk = text[start:end]
        chunks.append(chunk)
        start = end
    return chunks


def split_for_excel(text: str, max_chars: int = 32767) -> List[str]:
    if not text:
        return [""]
    return split_text(text, max_chars)


def sanitize_excel_value(val: Any) -> Any:
    if isinstance(val, str):
        return ILLEGAL_EXCEL_CHARS_RE.sub("", val)
    return val


def sanitize_excel_df(df: pd.DataFrame) -> pd.DataFrame:
    try:
        obj_cols = df.select_dtypes(include=["object"]).columns
        for col in obj_cols:
            df[col] = df[col].apply(sanitize_excel_value)
    except Exception:
        pass
    return df


def ocr_image_bytes(image_bytes: bytes) -> str:
    if pytesseract is None or Image is None:
        return ""
    try:
        with Image.open(io.BytesIO(image_bytes)) as img:
            return pytesseract.image_to_string(img) or ""
    except Exception:
        return ""


def extract_text_docx(path: Path, ocr_images: bool = False) -> str:
    if docx is None:
        raise RuntimeError("python-docx is not installed")
    doc = docx.Document(str(path))
    parts: List[str] = []
    for p in doc.paragraphs:
        if p.text:
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text or "" for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    if ocr_images:
        try:
            ocr_texts = []
            for shape in doc.inline_shapes:
                try:
                    r_id = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
                    image_part = doc.part.related_parts.get(r_id)
                    if image_part is None:
                        continue
                    txt = ocr_image_bytes(image_part.blob)
                    if txt.strip():
                        ocr_texts.append(txt.strip())
                except Exception:
                    continue
            if ocr_texts:
                logging.info("OCR extracted text from %d images in %s", len(ocr_texts), path)
                parts.append("\n".join(ocr_texts))
        except Exception:
            pass
    return "\n".join(parts)


def convert_doc_to_docx(path: Path) -> Optional[Path]:
    """Convert legacy .doc to .docx using LibreOffice (soffice)."""
    soffice = shutil.which("soffice") or shutil.which("soffice.exe")
    if not soffice:
        logging.warning("LibreOffice (soffice) not found; cannot convert %s", path)
        return None
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            out_dir = Path(tmpdir)
            cmd = [
                soffice,
                "--headless",
                "--convert-to",
                "docx",
                "--outdir",
                str(out_dir),
                str(path),
            ]
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            candidates = list(out_dir.glob("*.docx"))
            if not candidates:
                return None
            return candidates[0]
    except Exception as exc:
        logging.exception("Failed to convert %s: %s", path, exc)
        return None


def extract_text_pptx(path: Path, ocr_images: bool = False) -> str:
    if pptx is None:
        raise RuntimeError("python-pptx is not installed")
    pres = pptx.Presentation(str(path))
    parts: List[str] = []
    for slide in pres.slides:
        for shape in slide.shapes:
            text = ""
            if hasattr(shape, "text"):
                text = shape.text or ""
            elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                text = shape.text_frame.text or ""
            if text.strip():
                parts.append(text.strip())
            if ocr_images and MSO_SHAPE_TYPE is not None:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_blob = shape.image.blob
                        txt = ocr_image_bytes(img_blob)
                        if txt.strip():
                            parts.append(txt.strip())
                    except Exception:
                        pass
    return "\n".join(parts)


def extract_text_xlsx(path: Path) -> str:
    if openpyxl is None:
        raise RuntimeError("openpyxl is not installed")
    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    parts: List[str] = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            row_vals = [str(v) for v in row if v not in (None, "")]
            if row_vals:
                parts.append("\t".join(row_vals))
    return "\n".join(parts)


def extract_text_pdf(path: Path, ocrmypdf_enabled: bool) -> Tuple[str, List[str], str]:
    if pdfplumber is None:
        raise RuntimeError("pdfplumber is not installed")
    text, page_texts = extract_pdf_text_with_pdfplumber(path)
    logging.info("PDF text extracted (chars=%d) for %s", len(text.strip()), path)
    if len(text.strip()) >= MIN_EXTRACTED_CHARS:
        return text, page_texts, "ok"

    # OCRmyPDF default (if available)
    if ocrmypdf_enabled and ocrmypdf is not None:
        logging.info("OCR triggered (OCRmyPDF) for %s", path)
        ocr_text, ocr_pages, ocr_status = ocrmypdf_ocr_pdf(path, force_ocr=True)
        if len(ocr_text.strip()) >= MIN_EXTRACTED_CHARS:
            return ocr_text, ocr_pages, "ocrmypdf_used"
        # If OCRmyPDF failed or produced no text, try Tesseract fallback
        logging.info("OCR triggered (Tesseract fallback) for %s", path)
        ocr_texts, status = ocr_pdf(path)
        if ocr_texts:
            return "\n".join(ocr_texts), ocr_texts, "ocrmypdf_failed_then_ocr_used"
        return text, page_texts, ocr_status

    # Fallback OCR (Tesseract)
    if ocrmypdf_enabled:
        logging.info("OCR triggered (Tesseract fallback; OCRmyPDF unavailable) for %s", path)
    else:
        logging.info("OCR triggered (Tesseract fallback; OCRmyPDF disabled) for %s", path)
    ocr_texts, status = ocr_pdf(path)
    if ocr_texts:
        return "\n".join(ocr_texts), ocr_texts, status
    return text, page_texts, status


def extract_pdf_text_with_pdfplumber(path: Path) -> Tuple[str, List[str]]:
    parts: List[str] = []
    page_texts: List[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            page_texts.append(text)
            if text.strip():
                parts.append(text)
    return "\n".join(parts), page_texts


def ocrmypdf_ocr_pdf(path: Path, force_ocr: bool = False) -> Tuple[str, List[str], str]:
    if ocrmypdf is None:
        return "", [], "no_text_ocrmypdf_unavailable"
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            out_pdf = Path(tmpdir) / "ocr.pdf"
            kwargs = {
                "skip_text": not force_ocr,
                "output_type": "pdf",
                "progress_bar": False,
            }
            if force_ocr:
                kwargs["force_ocr"] = True
            try:
                ocrmypdf.ocr(str(path), str(out_pdf), **kwargs)
            except TypeError:
                # Older ocrmypdf versions may not support force_ocr
                kwargs.pop("force_ocr", None)
                ocrmypdf.ocr(str(path), str(out_pdf), **kwargs)
            text, page_texts = extract_pdf_text_with_pdfplumber(out_pdf)
            if len(text.strip()) >= MIN_EXTRACTED_CHARS:
                return text, page_texts, "ocrmypdf_used"
            return text, page_texts, "no_text_ocrmypdf"
    except Exception:
        return "", [], "no_text_ocrmypdf_failed"


def ocr_pdf(path: Path) -> Tuple[List[str], str]:
    if pytesseract is None or convert_from_path is None:
        return [], "no_text_ocr_unavailable"
    try:
        images = convert_from_path(str(path))
    except Exception:
        return [], "no_text_ocr_failed"
    texts: List[str] = []
    for img in images:
        try:
            txt = pytesseract.image_to_string(img)
        except Exception:
            txt = ""
        texts.append(txt or "")
    combined = "\n".join(texts).strip()
    if len(combined) >= MIN_EXTRACTED_CHARS:
        return texts, "ocr_used"
    return texts, "no_text"


def is_heading(line: str) -> bool:
    line = (line or "").strip()
    if not line:
        return False
    if len(line) > 120:
        return False
    if len(line) < 6:
        return False
    if re.match(r"^(article|section)\s+\d+", line, re.IGNORECASE):
        return True
    if re.match(r"^\d+(\.|\))\s+\S+", line):
        return True
    if line.isupper() and len(line) >= 6:
        return True
    # Title Case heuristic
    words = line.split()
    if len(words) >= 3 and sum(1 for w in words if w[:1].isupper()) / len(words) > 0.6:
        return True
    return False


def split_pdf_into_articles(page_texts: List[str]) -> List[Tuple[str, str]]:
    lines: List[str] = []
    for page in page_texts:
        for line in (page or "").splitlines():
            cleaned = line.strip()
            if cleaned:
                lines.append(cleaned)
    if not lines:
        return []

    indices: List[int] = []
    titles: List[str] = []
    for i, line in enumerate(lines):
        if is_heading(line):
            if indices and i - indices[-1] < MIN_HEADING_GAP_LINES:
                continue
            indices.append(i)
            titles.append(line)

    if not indices:
        return [("Article 1", "\n".join(lines))]

    articles: List[Tuple[str, str]] = []
    for idx, start in enumerate(indices):
        end = indices[idx + 1] if idx + 1 < len(indices) else len(lines)
        title = titles[idx] if idx < len(titles) else f"Article {idx + 1}"
        body = "\n".join(lines[start + 1 : end]).strip()
        if not body:
            body = "\n".join(lines[start:end]).strip()
        articles.append((title, body))

    # Merge very small sections into neighbors to avoid noisy article splits
    merged: List[Tuple[str, str]] = []
    for title, body in articles:
        if not merged:
            merged.append((title, body))
            continue
        if len(body) < MIN_ARTICLE_BODY_CHARS:
            prev_title, prev_body = merged[-1]
            merged[-1] = (prev_title, (prev_body + "\n" + body).strip())
        else:
            merged.append((title, body))

    # Drop articles still too small by merging into previous
    final: List[Tuple[str, str]] = []
    for title, body in merged:
        if not final:
            final.append((title, body))
            continue
        if len(body) < MIN_ARTICLE_BODY_CHARS:
            prev_title, prev_body = final[-1]
            final[-1] = (prev_title, (prev_body + "\n" + body).strip())
        else:
            final.append((title, body))

    # If everything merged into a tiny blob, treat as a single article
    if final and len(final) == 1 and len(final[0][1]) < MIN_ARTICLE_BODY_CHARS:
        return [("Article 1", "\n".join(lines))]
    return final


def call_azure_chat(cfg: AzureConfig, messages: List[Dict[str, str]]) -> str:
    url = build_url(cfg.chat_base_url, cfg.chat_path, cfg.chat_deployment)
    headers = {cfg.api_key_header: cfg.chat_api_key, "Content-Type": "application/json"}
    payload: Dict[str, Any] = {
        "messages": messages,
        "temperature": 0.2,
    }
    if cfg.include_model_in_body:
        payload["model"] = cfg.chat_deployment

    params = {"api-version": cfg.api_version}

    delay = api_delay_sec()
    for attempt in range(5):
        if delay:
            time.sleep(delay)
        resp = requests.post(url, headers=headers, params=params, json=payload, timeout=120)
        if resp.status_code in (429, 500, 502, 503, 504):
            time.sleep(1.5 * (2 ** attempt))
            continue
        if resp.status_code >= 400:
            raise RuntimeError(f"Chat API error {resp.status_code}: {resp.text}")
        data = resp.json()
        choices = data.get("choices") or []
        if not choices:
            raise RuntimeError("Chat API returned no choices")
        message = choices[0].get("message", {})
        content = message.get("content", "") or ""
        in_chars = sum(len(m.get("content", "") or "") for m in messages)
        add_chat_usage(in_chars, len(content))
        return content
    raise RuntimeError("Chat API failed after retries")


def call_azure_embeddings(cfg: AzureConfig, text: str) -> List[float]:
    url = build_url(cfg.embeddings_base_url, cfg.embeddings_path, cfg.embeddings_deployment)
    headers = {cfg.api_key_header: cfg.embeddings_api_key, "Content-Type": "application/json"}
    payload: Dict[str, Any] = {"input": text}
    if cfg.include_model_in_body:
        payload["model"] = cfg.embeddings_deployment
    params = {"api-version": cfg.api_version}
    delay = api_delay_sec()
    for attempt in range(5):
        if delay:
            time.sleep(delay)
        resp = requests.post(url, headers=headers, params=params, json=payload, timeout=120)
        if resp.status_code in (429, 500, 502, 503, 504):
            time.sleep(1.5 * (2 ** attempt))
            continue
        if resp.status_code >= 400:
            raise RuntimeError(f"Embeddings API error {resp.status_code}: {resp.text}")
        data = resp.json()
        data_list = data.get("data") or []
        if not data_list:
            raise RuntimeError("Embeddings API returned no data")
        add_embed_usage(len(text))
        return data_list[0].get("embedding")
    raise RuntimeError("Embeddings API failed after retries")


def extract_json(text: str) -> Dict[str, Any]:
    text = text.strip()
    if text.startswith("{"):
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            pass
    # Try to find JSON block
    match = re.search(r"\{.*\}", text, re.DOTALL)
    if match:
        return json.loads(match.group(0))
    raise ValueError("Failed to parse JSON")


def summarize_document(cfg: AzureConfig, text: str, categories: List[str]) -> Dict[str, Any]:
    categories_list = categories + (["Other"] if "Other" not in categories else [])
    if UNREADABLE_CATEGORY not in categories_list:
        categories_list.append(UNREADABLE_CATEGORY)
    if len(text) <= MAX_CHARS_PER_CHUNK:
        return summarize_with_model(cfg, text, categories_list)

    chunk_summaries: List[str] = []
    for chunk in split_text(text, MAX_CHARS_PER_CHUNK):
        chunk_summary = summarize_chunk(cfg, chunk)
        chunk_summaries.append(chunk_summary)
    combined = "\n".join(chunk_summaries)
    return summarize_with_model(cfg, combined, categories_list)


def summarize_chunk(cfg: AzureConfig, text: str) -> str:
    messages = [
        {"role": "system", "content": "You are a precise summarizer."},
        {
            "role": "user",
            "content": (
                "Summarize this chunk in 5-8 bullet points. Only output bullet points.\n\n"
                f"Chunk:\n{text}"
            ),
        },
    ]
    return call_azure_chat(cfg, messages)


def summarize_with_model(cfg: AzureConfig, text: str, categories: List[str]) -> Dict[str, Any]:
    categories_str = ", ".join(categories)
    messages = [
        {"role": "system", "content": "You are an expert analyst. Output JSON only."},
        {
            "role": "user",
            "content": (
                "Given the document text, produce JSON with keys: "
                "long_summary (5-7 sentences), short_summary (1-2 sentences), "
                "category (one of the provided categories), tags (array of strings). "
                "Tags can be as many as needed but should be specific and not redundant. "
                "If multiple categories could apply, prefer the most specific product/application "
                "category over broad process or issue buckets (e.g., prefer 'SeqStudio' over "
                "'Troubleshooting' when both fit). "
                f"Categories: {categories_str}.\n\n"
                f"Document:\n{text}"
            ),
        },
    ]
    content = call_azure_chat(cfg, messages)
    return extract_json(content)


def summarize_article(cfg: AzureConfig, text: str) -> str:
    text = text[:MAX_ARTICLE_CHARS]
    messages = [
        {"role": "system", "content": "You are a precise technical writer."},
        {
            "role": "user",
            "content": (
                "Write a condensed article that preserves all key facts, findings, and conclusions. "
                "Length should be as long as needed to capture every important point (do not force "
                "a fixed sentence count). No bullets. Do not start with phrases like "
                "'The article is about' or 'This article'. Write directly.\n\n"
                f"Article:\n{text}"
            ),
        },
    ]
    return call_azure_chat(cfg, messages).strip()


def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    if a.size == 0 or b.size == 0:
        return 0.0
    denom = np.linalg.norm(a) * np.linalg.norm(b)
    if denom == 0:
        return 0.0
    return float(np.dot(a, b) / denom)


def detect_duplicates(
    items: List[Tuple[str, str, Optional[np.ndarray]]], threshold: float
) -> Tuple[Dict[str, str], Dict[str, float], Dict[str, str]]:
    """
    items: list of (item_id, hash, embedding)
    returns: duplicate_of, duplicate_score, duplicate_group_id
    """
    duplicate_of: Dict[str, str] = {}
    duplicate_score: Dict[str, float] = {}
    duplicate_group_id: Dict[str, str] = {}

    hash_to_primary: Dict[str, str] = {}
    primary_embeddings: Dict[str, np.ndarray] = {}
    group_counter = 1

    for item_id, hsh, emb in items:
        if hsh in hash_to_primary:
            primary = hash_to_primary[hsh]
            duplicate_of[item_id] = primary
            duplicate_score[item_id] = 1.0
            group = duplicate_group_id.get(primary)
            if not group:
                group = f"DUP-{group_counter:04d}"
                group_counter += 1
                duplicate_group_id[primary] = group
            duplicate_group_id[item_id] = group
            continue

        # Near-duplicate check
        best_primary = None
        best_score = 0.0
        if emb is not None and primary_embeddings:
            for pid, pvec in primary_embeddings.items():
                score = cosine_similarity(emb, pvec)
                if score > best_score:
                    best_score = score
                    best_primary = pid
        if best_primary is not None and best_score >= threshold:
            duplicate_of[item_id] = best_primary
            duplicate_score[item_id] = best_score
            group = duplicate_group_id.get(best_primary)
            if not group:
                group = f"DUP-{group_counter:04d}"
                group_counter += 1
                duplicate_group_id[best_primary] = group
            duplicate_group_id[item_id] = group
        else:
            hash_to_primary[hsh] = item_id
            if emb is not None:
                primary_embeddings[item_id] = emb

    return duplicate_of, duplicate_score, duplicate_group_id


def get_categories_gui(
    app_config: Dict[str, List[str]],
    config_path: Path,
) -> Tuple[List[str], Optional[str]]:
    if tk is None:
        raise RuntimeError("tkinter is not available")

    root = tk.Tk()
    root.title("DocAtlas")
    root.geometry("720x560")
    apply_theme(root)

    selected_app = tk.StringVar(value="")

    container = tk.Frame(root, bg=THEME["bg"])
    container.pack(fill=tk.BOTH, expand=True, padx=16, pady=16)

    header = tk.Label(container, text="DocAtlas", bg=THEME["bg"], fg=THEME["fg"], font=FONT_HEADER)
    header.pack(anchor="w", pady=(0, 6))
    logo = tk.Label(container, text="— DocAtlas —", bg=THEME["bg"], fg=THEME["muted"], font=FONT_SMALL)
    logo.pack(anchor="w", pady=(0, 6))
    sub = tk.Label(
        container,
        text="Select Application & Categories",
        bg=THEME["bg"],
        fg=THEME["muted"],
        font=FONT_SMALL,
    )
    sub.pack(anchor="w", pady=(0, 12))

    btn_frame = tk.Frame(container, bg=THEME["bg"])
    btn_frame.pack(side=tk.BOTTOM, pady=12, fill=tk.X)

    panel = tk.Frame(container, bg=THEME["panel"], highlightbackground=THEME["border"], highlightthickness=1)
    panel.pack(fill=tk.BOTH, expand=True)

    if app_config:
        label_app = tk.Label(panel, text="Application", bg=THEME["panel"], fg=THEME["fg"], font=FONT_LABEL)
        label_app.pack(anchor="w", padx=12, pady=(12, 4))
        options = ["(Custom)"] + sorted(app_config.keys())
        dropdown = tk.OptionMenu(panel, selected_app, *options)
        dropdown.config(bg=THEME["text_bg"], fg=THEME["fg"], highlightthickness=0, font=FONT_BASE)
        try:
            dropdown["menu"].config(bg=THEME["text_bg"], fg=THEME["fg"], font=FONT_BASE)
        except Exception:
            pass
        dropdown.pack(anchor="w", padx=12, pady=(0, 8))
        selected_app.set(options[0])

    label = tk.Label(panel, text="Categories (one per line)", bg=THEME["panel"], fg=THEME["fg"], font=FONT_LABEL)
    label.pack(anchor="w", padx=12, pady=(8, 4))

    text_widget = tk.Text(
        panel,
        height=16,
        width=78,
        bg=THEME["text_bg"],
        fg=THEME["fg"],
        insertbackground=THEME["fg"],
        highlightbackground=THEME["border"],
        font=FONT_BASE,
    )
    text_widget.pack(padx=12, pady=(0, 12), fill=tk.BOTH, expand=True)

    result: List[str] = []
    result_app: Optional[str] = None

    def apply_app_categories(*_args: Any) -> None:
        app = selected_app.get()
        if app and app in app_config:
            text_widget.delete("1.0", tk.END)
            text_widget.insert(tk.END, "\n".join(app_config[app]))

    if app_config:
        selected_app.trace_add("write", apply_app_categories)

    def on_ok() -> None:
        content = text_widget.get("1.0", tk.END)
        lines = [line.strip() for line in content.splitlines() if line.strip()]
        if not lines:
            messagebox.showerror("Error", "Please enter at least one category.")
            return
        nonlocal result
        result = lines
        nonlocal result_app
        app = selected_app.get().strip()
        if app and app != "(Custom)" and app in app_config:
            result_app = app
        root.destroy()

    def on_cancel() -> None:
        root.destroy()

    if app_config:
        tk.Button(
            btn_frame,
            text="Edit Apps",
            command=lambda: edit_applications_gui(config_path, app_config, root),
            width=12,
            bg=THEME["btn_bg"],
            fg=THEME["btn_fg"],
            relief=tk.FLAT,
            borderwidth=1,
            font=FONT_BUTTON,
        ).pack(side=tk.RIGHT, padx=8, ipady=6)
    tk.Button(btn_frame, text="Cancel", command=on_cancel, width=12, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.RAISED, borderwidth=1, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8, ipady=6)
    tk.Button(btn_frame, text="OK", command=on_ok, width=12, bg=THEME["accent"], fg="#ffffff", relief=tk.RAISED, borderwidth=1, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8, ipady=6)

    root.mainloop()
    return result, result_app


def pick_directories_gui() -> Tuple[Path, Path]:
    if tk is None:
        raise RuntimeError("tkinter is not available")

    root = tk.Tk()
    root.withdraw()
    input_dir = filedialog.askdirectory(title="Select Input Folder")
    if not input_dir:
        raise RuntimeError("No input folder selected")

    output_dir = filedialog.askdirectory(
        title="Select Output Folder (Cancel = use input folder)",
        initialdir=input_dir,
    )
    if not output_dir:
        output_dir = input_dir

    root.destroy()
    return Path(input_dir), Path(output_dir)


def get_ocrmypdf_gui() -> bool:
    if tk is None:
        return True
    root = tk.Tk()
    root.title("DocAtlas - OCR Options")
    root.geometry("520x260")
    apply_theme(root)

    var = tk.BooleanVar(value=True)
    container = tk.Frame(root, bg=THEME["bg"])
    container.pack(fill=tk.BOTH, expand=True, padx=16, pady=16)

    label = tk.Label(container, text="OCR Settings", bg=THEME["bg"], fg=THEME["fg"], font=FONT_HEADER)
    label.pack(anchor="w", pady=(0, 6))
    sub = tk.Label(container, text="OCR runs only when PDFs have little or no text.", bg=THEME["bg"], fg=THEME["muted"], font=FONT_SMALL)
    sub.pack(anchor="w", pady=(0, 12))

    chk = tk.Checkbutton(container, text="Use OCRmyPDF (recommended)", variable=var, bg=THEME["bg"], fg=THEME["fg"], font=FONT_LABEL)
    chk.pack(anchor="w", pady=6)

    result: List[bool] = []

    def on_ok() -> None:
        result.append(bool(var.get()))
        root.destroy()

    btn_frame = tk.Frame(container, bg=THEME["bg"])
    btn_frame.pack(pady=12, anchor="e", fill=tk.X)

    def on_test() -> None:
        missing = warn_missing_ocr_deps(bool(var.get()))
        if not missing and messagebox is not None:
            messagebox.showinfo("OCR Dependencies", "All OCR dependencies found.")

    tk.Button(btn_frame, text="Test OCR", command=on_test, width=12, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.FLAT, height=1, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8, ipady=4)
    tk.Button(btn_frame, text="OK", command=on_ok, width=12, bg=THEME["accent"], fg="#ffffff", relief=tk.FLAT, height=1, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8, ipady=4)
    root.mainloop()
    return result[0] if result else True


def get_embeddings_source_gui() -> Tuple[str, bool]:
    if tk is None:
        return DEFAULT_EMBEDDINGS_SOURCE, True
    root = tk.Tk()
    root.title("DocAtlas - Embeddings Source")
    root.geometry("620x340")
    apply_theme(root)

    var = tk.StringVar(value=DEFAULT_EMBEDDINGS_SOURCE)
    append_var = tk.BooleanVar(value=True)
    container = tk.Frame(root, bg=THEME["bg"])
    container.pack(fill=tk.BOTH, expand=True, padx=16, pady=16)

    label = tk.Label(container, text="Embeddings Source", bg=THEME["bg"], fg=THEME["fg"], font=FONT_HEADER)
    label.pack(anchor="w", pady=(0, 6))
    sub = tk.Label(
        container,
        text="Choose what text is embedded for duplicate detection.",
        bg=THEME["bg"],
        fg=THEME["muted"],
        font=FONT_SMALL,
    )
    sub.pack(anchor="w", pady=(0, 12))

    tk.Radiobutton(
        container,
        text="Long summary (lower cost)",
        variable=var,
        value="summary",
        bg=THEME["bg"],
        fg=THEME["fg"],
        selectcolor=THEME["panel"],
        font=FONT_LABEL,
    ).pack(anchor="w", pady=4)
    tk.Radiobutton(
        container,
        text="Full text (recommended, stricter)",
        variable=var,
        value="full_text",
        bg=THEME["bg"],
        fg=THEME["fg"],
        selectcolor=THEME["panel"],
        font=FONT_LABEL,
    ).pack(anchor="w", pady=4)
    tk.Radiobutton(
        container,
        text="Disable embeddings (hash-only duplicates)",
        variable=var,
        value=EMBEDDINGS_SOURCE_NONE,
        bg=THEME["bg"],
        fg=THEME["fg"],
        selectcolor=THEME["panel"],
        font=FONT_LABEL,
    ).pack(anchor="w", pady=4)

    chk = tk.Checkbutton(
        container,
        text="Append to existing Excel (recommended)",
        variable=append_var,
        bg=THEME["bg"],
        fg=THEME["fg"],
        font=FONT_LABEL,
    )
    chk.pack(anchor="w", pady=(8, 4))

    result: List[Tuple[str, bool]] = []

    def on_ok() -> None:
        result.append((var.get(), bool(append_var.get())))
        root.destroy()

    btn_frame = tk.Frame(container, bg=THEME["bg"])
    btn_frame.pack(pady=12, anchor="e", fill=tk.X)
    tk.Button(btn_frame, text="OK", command=on_ok, width=12, bg=THEME["accent"], fg="#ffffff", relief=tk.FLAT, height=1, font=FONT_BUTTON).pack(
        side=tk.RIGHT, padx=8, ipady=4
    )
    root.mainloop()
    return result[0] if result else (DEFAULT_EMBEDDINGS_SOURCE, True)


def get_run_mode_gui() -> Optional[bool]:
    if tk is None:
        return True
    root = tk.Tk()
    root.title("DocAtlas - Run Mode")
    root.geometry("560x300")
    apply_theme(root)

    var = tk.StringVar(value="charter")
    container = tk.Frame(root, bg=THEME["bg"])
    container.pack(fill=tk.BOTH, expand=True, padx=16, pady=16)

    label = tk.Label(container, text="Run Mode", bg=THEME["bg"], fg=THEME["fg"], font=FONT_HEADER)
    label.pack(anchor="w", pady=(0, 6))
    sub = tk.Label(
        container,
        text="Choose how to run this job.",
        bg=THEME["bg"],
        fg=THEME["muted"],
        font=FONT_SMALL,
    )
    sub.pack(anchor="w", pady=(0, 12))

    tk.Radiobutton(
        container,
        text="Charter Mode (summaries, tags, duplicates; no file moves)",
        variable=var,
        value="charter",
        bg=THEME["bg"],
        fg=THEME["fg"],
        selectcolor=THEME["panel"],
        font=FONT_LABEL,
    ).pack(anchor="w", pady=4)
    tk.Radiobutton(
        container,
        text="Atlas Mode (summaries, tags, duplicates; move files)",
        variable=var,
        value="atlas",
        bg=THEME["bg"],
        fg=THEME["fg"],
        selectcolor=THEME["panel"],
        font=FONT_LABEL,
    ).pack(anchor="w", pady=4)

    result: List[Optional[bool]] = []

    def on_ok() -> None:
        result.append(var.get() == "charter")
        root.destroy()

    def on_close() -> None:
        result.append(None)
        root.destroy()

    root.protocol("WM_DELETE_WINDOW", on_close)

    btn_frame = tk.Frame(container, bg=THEME["bg"])
    btn_frame.pack(pady=12, anchor="e", fill=tk.X)
    tk.Button(btn_frame, text="OK", command=on_ok, width=12, bg=THEME["accent"], fg="#ffffff", relief=tk.FLAT, height=1, font=FONT_BUTTON).pack(
        side=tk.RIGHT, padx=8, ipady=4
    )
    root.mainloop()
    return result[0] if result else True


def edit_applications_gui(config_path: Path, app_config: Dict[str, List[str]], parent: tk.Tk) -> None:
    if tk is None:
        return
    win = tk.Toplevel(parent)
    win.title("DocAtlas - Edit Applications")
    win.geometry("820x560")
    apply_theme(win)

    apps = dict(app_config)

    left_frame = tk.Frame(win, bg=THEME["panel"], highlightbackground=THEME["border"], highlightthickness=1)
    left_frame.pack(side=tk.LEFT, fill=tk.Y, padx=10, pady=10)

    right_frame = tk.Frame(win, bg=THEME["bg"])
    right_frame.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True, padx=10, pady=10)

    listbox = tk.Listbox(
        left_frame,
        height=20,
        width=25,
        bg=THEME["text_bg"],
        fg=THEME["fg"],
        selectbackground=THEME["accent"],
        selectforeground="#ffffff",
        font=FONT_BASE,
    )
    listbox.pack(pady=5)

    def refresh_list() -> None:
        listbox.delete(0, tk.END)
        for name in sorted(apps.keys()):
            listbox.insert(tk.END, name)

    def load_selected(event: Any = None) -> None:
        selection = listbox.curselection()
        if not selection:
            return
        name = listbox.get(selection[0])
        name_entry.delete(0, tk.END)
        name_entry.insert(0, name)
        cat_text.delete("1.0", tk.END)
        cat_text.insert(tk.END, "\n".join(apps.get(name, [])))

    def add_new() -> None:
        name_entry.delete(0, tk.END)
        cat_text.delete("1.0", tk.END)

    def save_current() -> None:
        name = name_entry.get().strip()
        if not name:
            if messagebox:
                messagebox.showerror("Error", "Application name is required.")
            return
        cats = [c.strip() for c in cat_text.get("1.0", tk.END).splitlines() if c.strip()]
        apps[name] = cats
        refresh_list()

    def delete_current() -> None:
        name = name_entry.get().strip()
        if not name:
            return
        if name in apps:
            del apps[name]
        refresh_list()
        name_entry.delete(0, tk.END)
        cat_text.delete("1.0", tk.END)

    def save_all() -> None:
        save_current()
        save_app_config(config_path, apps)
        app_config.clear()
        app_config.update(apps)
        if messagebox:
            messagebox.showinfo("Saved", f"Saved to {config_path}")

    def save_and_close() -> None:
        save_all()
        win.destroy()

    listbox.bind("<<ListboxSelect>>", load_selected)

    tk.Label(right_frame, text="Application Name", bg=THEME["bg"], fg=THEME["fg"], font=FONT_LABEL).pack(anchor="w")
    name_entry = tk.Entry(right_frame, width=40, bg=THEME["text_bg"], fg=THEME["fg"], font=FONT_BASE)
    name_entry.pack(fill=tk.X, pady=4)

    tk.Label(right_frame, text="Categories (one per line)", bg=THEME["bg"], fg=THEME["fg"], font=FONT_LABEL).pack(anchor="w")
    cat_text = tk.Text(
        right_frame,
        height=15,
        bg=THEME["text_bg"],
        fg=THEME["fg"],
        insertbackground=THEME["fg"],
        highlightbackground=THEME["border"],
        font=FONT_BASE,
    )
    cat_text.pack(fill=tk.BOTH, expand=True)

    btn_row = tk.Frame(right_frame, bg=THEME["bg"])
    btn_row.pack(pady=8, fill=tk.X)
    tk.Button(btn_row, text="Save & Close", command=save_and_close, width=14, bg=THEME["accent"], fg="#ffffff", relief=tk.FLAT, height=1, padx=10, pady=6, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8)
    tk.Button(btn_row, text="Save All", command=save_all, width=12, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.FLAT, height=1, padx=10, pady=6, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8)
    tk.Button(btn_row, text="Delete", command=delete_current, width=10, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.FLAT, height=1, padx=10, pady=6, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8)
    tk.Button(btn_row, text="Save App", command=save_current, width=10, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.FLAT, height=1, padx=10, pady=6, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8)
    tk.Button(btn_row, text="New", command=add_new, width=10, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.FLAT, height=1, padx=10, pady=6, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8)

    refresh_list()


def process_file(path: Path, ocrmypdf_enabled: bool) -> Tuple[str, List[Tuple[str, str]], str]:
    ext = path.suffix.lower()
    try:
        if path.stat().st_size == 0:
            return "", [], "no_text_empty"
    except Exception:
        pass
    if ext == ".docx":
        text = extract_text_docx(path, ocrmypdf_enabled)
        status = "ok" if len(text.strip()) >= MIN_EXTRACTED_CHARS else "no_text"
        return text, [], status
    if ext == ".doc":
        docx_path = convert_doc_to_docx(path)
        if docx_path is None:
            return "", [], "no_text_doc_convert_failed"
        text = extract_text_docx(docx_path, ocrmypdf_enabled)
        status = "ok" if len(text.strip()) >= MIN_EXTRACTED_CHARS else "no_text"
        return text, [], status
    if ext == ".pptx":
        text = extract_text_pptx(path, ocrmypdf_enabled)
        status = "ok" if len(text.strip()) >= MIN_EXTRACTED_CHARS else "no_text"
        return text, [], status
    if ext == ".xlsx":
        text = extract_text_xlsx(path)
        status = "ok" if len(text.strip()) >= MIN_EXTRACTED_CHARS else "no_text"
        return text, [], status
    if ext == ".pdf":
        text, pages, status = extract_text_pdf(path, ocrmypdf_enabled)
        articles = split_pdf_into_articles(pages)
        return text, articles, status
    raise ValueError(f"Unsupported file type: {ext}")


def write_excels(
    out_dir: Path,
    docs: List[DocRecord],
    articles: List[ArticleRecord],
    full_text_rows: List[Dict[str, Any]],
    app_name: Optional[str],
    append_excel: bool,
) -> Tuple[Path, Path]:
    app_slug = sanitize_folder(app_name or "uncategorized")
    peers_path = out_dir / f"{app_slug}__docatlas_summaries.xlsx"
    full_text_path = out_dir / f"{app_slug}__docatlas_full_text.xlsx"

    existing_docs_df = None
    existing_articles_df = None
    existing_full_df = None
    existing_doc_keys: set[str] = set()

    if append_excel and peers_path.exists():
        try:
            existing_docs_df = pd.read_excel(peers_path, sheet_name="Documents")
            if "file_key" in existing_docs_df.columns:
                existing_doc_keys = set(existing_docs_df["file_key"].astype(str))
            elif "file_path" in existing_docs_df.columns:
                existing_doc_keys = set(existing_docs_df["file_path"].astype(str))
        except Exception:
            existing_docs_df = None
            existing_doc_keys = set()
        try:
            existing_articles_df = pd.read_excel(peers_path, sheet_name="Articles")
        except Exception:
            existing_articles_df = None

    if append_excel and full_text_path.exists():
        try:
            existing_full_df = pd.read_excel(full_text_path, sheet_name="FullText")
        except Exception:
            existing_full_df = None

    docs_rows = []
    new_doc_ids: set[str] = set()
    for d in docs:
        key = d.file_key or d.file_path
        if append_excel and key in existing_doc_keys:
            continue
        new_doc_ids.add(d.doc_id)
        docs_rows.append(
            {
                "doc_id": d.doc_id,
                "file_key": d.file_key,
                "category": d.category,
                "tags": ", ".join(d.tags),
                "short_summary": d.short_summary,
                "long_summary": d.long_summary,
                "word_count": d.word_count,
                "char_count": d.char_count,
                "extraction_status": d.extraction_status,
                "review_flags": d.review_flags,
                "duplicate_of": d.duplicate_of,
                "duplicate_score": d.duplicate_score,
                "duplicate_group_id": d.duplicate_group_id,
                "moved_to": d.moved_to,
                "file_name": d.file_name,
                "file_path": d.file_path,
            }
        )

    articles_rows = []
    for a in articles:
        if append_excel and a.doc_id not in new_doc_ids:
            continue
        articles_rows.append(
            {
                "doc_id": a.doc_id,
                "file_key": a.file_key,
                "article_index": a.article_index,
                "article_title": a.article_title,
                "article_summary": a.article_summary,
                "duplicate_of": a.duplicate_of,
                "duplicate_score": a.duplicate_score,
                "duplicate_group_id": a.duplicate_group_id,
                "file_name": a.file_name,
                "file_path": a.file_path,
            }
        )

    docs_df = sanitize_excel_df(pd.DataFrame(docs_rows))
    articles_df = sanitize_excel_df(pd.DataFrame(articles_rows))
    if append_excel and existing_docs_df is not None:
        docs_df = pd.concat([existing_docs_df, docs_df], ignore_index=True)
    if append_excel and existing_articles_df is not None:
        articles_df = pd.concat([existing_articles_df, articles_df], ignore_index=True)

    with pd.ExcelWriter(peers_path, engine="openpyxl") as writer:
        docs_df.to_excel(writer, index=False, sheet_name="Documents")
        articles_df.to_excel(writer, index=False, sheet_name="Articles")

    # Expand full_text parts into separate columns to avoid Excel 32,767 char limit
    expanded_rows: List[Dict[str, Any]] = []
    max_parts = 1
    for row in full_text_rows:
        parts = row.pop("full_text_parts", [])
        max_parts = max(max_parts, len(parts))
        row["_parts"] = parts
        expanded_rows.append(row)

    for row in expanded_rows:
        parts = row.pop("_parts", [])
        for i in range(max_parts):
            key = f"full_text_part_{i+1}"
            row[key] = parts[i] if i < len(parts) else ""

    full_df = sanitize_excel_df(pd.DataFrame(expanded_rows))
    if append_excel and existing_doc_keys:
        if "file_key" in full_df.columns:
            full_df = full_df[~full_df["file_key"].astype(str).isin(existing_doc_keys)]
        elif "file_path" in full_df.columns:
            full_df = full_df[~full_df["file_path"].astype(str).isin(existing_doc_keys)]
    if append_excel and existing_full_df is not None:
        full_df = pd.concat([existing_full_df, full_df], ignore_index=True)

    with pd.ExcelWriter(full_text_path, engine="openpyxl") as writer:
        full_df.to_excel(writer, index=False, sheet_name="FullText")

    # Apply formatting: wrap summaries and widen columns
    try:
        from openpyxl import load_workbook
        from openpyxl.styles import Alignment
        from openpyxl.worksheet.table import Table, TableStyleInfo

        def format_sheet(path: Path, sheet_name: str, wrap_cols: List[str]) -> None:
            wb = load_workbook(path)
            ws = wb[sheet_name]
            header = {cell.value: cell.column for cell in ws[1]}
            wrap = Alignment(wrap_text=True, vertical="top")

            for col_name in wrap_cols:
                col_idx = header.get(col_name)
                if not col_idx:
                    continue
                for row in ws.iter_rows(min_row=2, min_col=col_idx, max_col=col_idx):
                    for cell in row:
                        cell.alignment = wrap
                # widen column
                col_letter = ws.cell(row=1, column=col_idx).column_letter
                ws.column_dimensions[col_letter].width = 60

            # Apply table style
            max_row = ws.max_row
            max_col = ws.max_column
            if max_row >= 2 and max_col >= 1:
                ref = f"A1:{ws.cell(row=max_row, column=max_col).coordinate}"
                table = Table(displayName=f"{sheet_name}Table", ref=ref)
                style = TableStyleInfo(
                    name="TableStyleMedium1",
                    showFirstColumn=False,
                    showLastColumn=False,
                    showRowStripes=True,
                    showColumnStripes=False,
                )
                table.tableStyleInfo = style
                ws.add_table(table)

            wb.save(path)

        format_sheet(peers_path, "Documents", ["short_summary", "long_summary"])
        format_sheet(full_text_path, "FullText", ["short_summary", "long_summary", "full_text"])
    except Exception:
        pass

    return peers_path, full_text_path


def write_summary_report(
    out_dir: Path,
    docs: List[DocRecord],
    articles: List[ArticleRecord],
    errors: Optional[List[Dict[str, str]]] = None,
    usage: Optional[Dict[str, int]] = None,
    total_files: Optional[int] = None,
    processed_files: Optional[int] = None,
    limit: Optional[int] = None,
) -> Path:
    report_path = out_dir / "summary_report.txt"
    total_docs = len(docs)
    total_articles = len(articles)
    dup_docs = sum(1 for d in docs if d.duplicate_of)
    dup_articles = sum(1 for a in articles if a.duplicate_of)

    categories: Dict[str, int] = {}
    for d in docs:
        categories[d.category] = categories.get(d.category, 0) + 1

    ext_counts: Dict[str, int] = {}
    for d in docs:
        ext_counts[d.file_ext] = ext_counts.get(d.file_ext, 0) + 1

    extraction_status: Dict[str, int] = {}
    for d in docs:
        extraction_status[d.extraction_status] = extraction_status.get(d.extraction_status, 0) + 1
    no_text_docs = [d for d in docs if d.extraction_status != "ok"]
    ocr_docs = [d for d in docs if d.extraction_status.startswith("ocr") or "ocr" in d.extraction_status]

    word_counts = [d.word_count for d in docs if d.word_count is not None]
    char_counts = [d.char_count for d in docs if d.char_count is not None]
    avg_words = int(sum(word_counts) / len(word_counts)) if word_counts else 0
    avg_chars = int(sum(char_counts) / len(char_counts)) if char_counts else 0
    longest = max(docs, key=lambda d: d.char_count or 0, default=None)
    shortest = min(docs, key=lambda d: d.char_count or 0, default=None)

    dup_group_sizes: Dict[str, int] = {}
    for d in docs:
        if d.duplicate_group_id:
            dup_group_sizes[d.duplicate_group_id] = dup_group_sizes.get(d.duplicate_group_id, 0) + 1
    dup_group_count = len(dup_group_sizes)
    avg_dup_group_size = int(sum(dup_group_sizes.values()) / dup_group_count) if dup_group_count else 0

    lines = []
    lines.append("Summary Report")
    lines.append("================")
    lines.append(f"Total documents: {total_docs}")
    lines.append(f"Total articles: {total_articles}")
    lines.append(f"Duplicate documents: {dup_docs}")
    lines.append(f"Duplicate articles: {dup_articles}")
    lines.append("")
    lines.append("Documents by Category:")
    for k in sorted(categories.keys()):
        pct = (categories[k] / total_docs * 100) if total_docs else 0
        lines.append(f"- {k}: {categories[k]} ({pct:.1f}%)")
    lines.append("")
    lines.append("Documents by File Type:")
    for k in sorted(ext_counts.keys()):
        pct = (ext_counts[k] / total_docs * 100) if total_docs else 0
        lines.append(f"- {k}: {ext_counts[k]} ({pct:.1f}%)")
    lines.append("")
    lines.append("Document Length (approx):")
    lines.append(f"- avg_words: {avg_words}")
    lines.append(f"- avg_chars: {avg_chars}")
    if longest:
        lines.append(f"- longest: {longest.file_name} ({longest.char_count} chars)")
    if shortest:
        lines.append(f"- shortest: {shortest.file_name} ({shortest.char_count} chars)")
    lines.append("")
    lines.append("OCR Usage:")
    lines.append(f"- ocr_used: {len(ocr_docs)}")
    lines.append("")
    lines.append("Duplicate Groups:")
    lines.append(f"- duplicate_group_count: {dup_group_count}")
    lines.append(f"- avg_duplicate_group_size: {avg_dup_group_size}")
    lines.append("")
    lines.append("Extraction Status:")
    for k in sorted(extraction_status.keys()):
        lines.append(f"- {k}: {extraction_status[k]}")
    if no_text_docs:
        lines.append("")
        lines.append("No-Text Files:")
        for d in no_text_docs[:50]:
            lines.append(f"- {d.file_name}")
        if len(no_text_docs) > 50:
            lines.append(f"- ... ({len(no_text_docs) - 50} more)")

    if errors:
        lines.append("")
        lines.append("Errors:")
        for e in errors[:50]:
            lines.append(f"- {e.get('stage','unknown')}: {e.get('file_name','')} | {e.get('error','')}")
        if len(errors) > 50:
            lines.append(f"- ... ({len(errors) - 50} more)")

    if usage:
        chat_in = usage.get("chat_in", 0)
        chat_out = usage.get("chat_out", 0)
        embed_in = usage.get("embed_in", 0)
        tokens_in = int(chat_in / 4)
        tokens_out = int(chat_out / 4)
        tokens_embed = int(embed_in / 4)
        total_tokens = tokens_in + tokens_out + tokens_embed
        lines.append("")
        lines.append("Token Estimates (approx):")
        lines.append(f"- chat_input_tokens: {tokens_in}")
        lines.append(f"- chat_output_tokens: {tokens_out}")
        lines.append(f"- embeddings_tokens: {tokens_embed}")
        lines.append(f"- total_tokens: {total_tokens}")

    report_path.write_text("\n".join(lines), encoding="utf-8")
    return report_path


def prompt_api_key_gui(title: str, label_text: str) -> Optional[str]:
    if tk is None:
        return None
    root = tk.Tk()
    root.title(title)
    root.geometry("560x240")
    apply_theme(root)

    container = tk.Frame(root, bg=THEME["bg"])
    container.pack(fill=tk.BOTH, expand=True, padx=16, pady=16)

    label = tk.Label(container, text=label_text, bg=THEME["bg"], fg=THEME["fg"], font=FONT_LABEL)
    label.pack(anchor="w", pady=(0, 8))

    entry = tk.Entry(container, show="*", width=60, bg=THEME["text_bg"], fg=THEME["fg"], insertbackground=THEME["fg"], font=FONT_BASE)
    entry.pack(fill=tk.X, pady=(0, 8))

    result: List[str] = []

    def on_ok() -> None:
        val = entry.get().strip()
        if not val:
            if messagebox:
                messagebox.showerror("Error", "API key is required.")
            return
        result.append(val)
        root.destroy()

    def on_cancel() -> None:
        root.destroy()

    btn_frame = tk.Frame(container, bg=THEME["bg"])
    btn_frame.pack(pady=12, anchor="e", fill=tk.X)
    tk.Button(btn_frame, text="Cancel", command=on_cancel, width=12, bg=THEME["btn_bg"], fg=THEME["btn_fg"], relief=tk.FLAT, height=1, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8, ipady=6)
    tk.Button(btn_frame, text="OK", command=on_ok, width=12, bg=THEME["accent"], fg="#ffffff", relief=tk.FLAT, height=1, font=FONT_BUTTON).pack(side=tk.RIGHT, padx=8, ipady=6)

    root.mainloop()
    return result[0] if result else None


def file_key(path: Path) -> str:
    st = path.stat()
    return f"{path}|{st.st_mtime_ns}|{st.st_size}"


def load_resume(out_dir: Path) -> Dict[str, Any]:
    resume_path = out_dir / RESUME_FILENAME
    if not resume_path.exists():
        return {"files": {}}
    try:
        with resume_path.open("r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {"files": {}}


def save_resume(out_dir: Path, resume: Dict[str, Any]) -> None:
    resume_path = out_dir / RESUME_FILENAME
    with resume_path.open("w", encoding="utf-8") as f:
        json.dump(resume, f)


def run_pipeline(
    input_dir: Path,
    output_dir: Path,
    categories: List[str],
    cfg: AzureConfig,
    dry_run: bool,
    use_resume: bool,
    ocrmypdf_enabled: bool,
    app_name: Optional[str],
    embeddings_source: str,
    append_excel: bool,
    limit: Optional[int] = None,
    no_move: bool = False,
    progress_cb: Optional[callable] = None,
) -> None:
    setup_logging(output_dir)
    logging.info("Starting pipeline")

    reset_usage()

    files = list_files(input_dir)
    total_files = len(files)
    input_stats = scan_input_stats(files)
    if limit is not None and limit > 0:
        files = files[:limit]
    processed_stats = scan_input_stats(files)
    if not files:
        logging.warning("No supported files found")
        return

    run_id = time.strftime("%Y%m%d%H%M%S")
    t0 = time.time()

    resume = load_resume(output_dir) if use_resume else {"files": {}}
    resume_files = resume.get("files", {})

    doc_items: List[Tuple[str, str, Optional[np.ndarray]]] = []
    article_items: List[Tuple[str, str, Optional[np.ndarray]]] = []
    docs: List[DocRecord] = []
    articles: List[ArticleRecord] = []

    doc_hashes: Dict[str, str] = {}
    article_hashes: Dict[str, str] = {}
    doc_id_to_key: Dict[str, str] = {}
    article_id_to_key: Dict[str, str] = {}
    article_id_to_idx: Dict[str, int] = {}

    raw_texts: Dict[str, str] = {}
    article_texts: Dict[str, str] = {}
    extraction_statuses: Dict[str, str] = {}
    errors: List[Dict[str, str]] = []

    iterable = tqdm(files, desc="Extracting") if tqdm else files
    for idx, path in enumerate(iterable, start=1):
        if progress_cb:
            progress_cb("Extracting", idx - 1, len(files))
        key = file_key(path)
        cached = resume_files.get(key)
        doc_id = f"{run_id}-DOC-{idx:05d}"
        if cached and cached.get("doc_id"):
            doc_id = cached["doc_id"]
        logging.info("Processing %s", path)
        doc_id_to_key[doc_id] = key
        if cached:
            text = cached.get("text", "")
            pdf_articles = cached.get("articles_raw", [])
            extraction_status = cached.get("extraction_status", "no_text")
            cached["doc_id"] = doc_id
        else:
            try:
                text, pdf_articles, extraction_status = process_file(path, ocrmypdf_enabled)
            except Exception as exc:
                logging.exception("Failed to extract %s: %s", path, exc)
                errors.append({"stage": "extract", "file_name": path.name, "file_path": str(path), "error": str(exc)})
                text = ""
                pdf_articles = []
                extraction_status = "no_text"
            resume_files[key] = {
                "doc_id": doc_id,
                "file_path": str(path),
                "file_name": path.name,
                "ext": path.suffix.lower(),
                "text": text,
                "articles_raw": pdf_articles,
                "extraction_status": extraction_status,
            }

        raw_texts[doc_id] = text
        extraction_statuses[doc_id] = resume_files.get(key, {}).get("extraction_status", extraction_status)
        normalized = normalize_text(text)
        hsh = hash_text(normalized)
        doc_hashes[doc_id] = hsh

        emb_vec: Optional[np.ndarray] = None
        if embeddings_source == "full_text":
            cached_emb = cached.get("doc_embedding") if cached else None
            if cached_emb is not None:
                emb_vec = np.array(cached_emb, dtype=np.float32)
            elif normalized and len(normalized) >= MIN_EMBEDDING_CHARS and not dry_run:
                try:
                    emb = call_azure_embeddings(cfg, normalized[:MAX_CHARS_PER_CHUNK])
                    emb_vec = np.array(emb, dtype=np.float32)
                    resume_files[key]["doc_embedding"] = emb
                except Exception as exc:
                    logging.exception("Embedding failed for %s: %s", path, exc)
                    errors.append({"stage": "embedding", "file_name": path.name, "file_path": str(path), "error": str(exc)})

        doc_items.append((doc_id, hsh, emb_vec))

        # Article handling (PDF only)
        for a_idx, (title, body) in enumerate(pdf_articles, start=1):
            article_id = f"{doc_id}-A{a_idx:03d}"
            article_texts[article_id] = body
            ahash = hash_text(normalize_text(body))
            article_hashes[article_id] = ahash
            article_id_to_key[article_id] = key
            article_id_to_idx[article_id] = a_idx
            aemb_vec: Optional[np.ndarray] = None
            if embeddings_source == "full_text":
                cached_aemb = None
                if cached and "article_embeddings" in cached:
                    cached_aemb = cached["article_embeddings"].get(str(a_idx))
                if cached_aemb is not None:
                    aemb_vec = np.array(cached_aemb, dtype=np.float32)
                elif body.strip() and len(body) >= MIN_EMBEDDING_CHARS and not dry_run:
                    try:
                        aemb = call_azure_embeddings(cfg, body[:MAX_CHARS_PER_CHUNK])
                        aemb_vec = np.array(aemb, dtype=np.float32)
                        resume_files[key].setdefault("article_embeddings", {})[str(a_idx)] = aemb
                    except Exception as exc:
                        logging.exception("Article embedding failed for %s: %s", path, exc)
                        errors.append({"stage": "article_embedding", "file_name": path.name, "file_path": str(path), "error": str(exc)})
            article_items.append((article_id, ahash, aemb_vec))

        if use_resume:
            resume["files"] = resume_files
            save_resume(output_dir, resume)
        if progress_cb:
            progress_cb("Extracting", idx, len(files))

    if embeddings_source == "full_text":
        doc_dup_of, doc_dup_score, doc_dup_group = detect_duplicates(doc_items, DUPLICATE_THRESHOLD)
        art_dup_of, art_dup_score, art_dup_group = detect_duplicates(article_items, DUPLICATE_THRESHOLD)
    elif embeddings_source == EMBEDDINGS_SOURCE_NONE:
        doc_dup_of, doc_dup_score, doc_dup_group = detect_duplicates(doc_items, DUPLICATE_THRESHOLD)
        art_dup_of, art_dup_score, art_dup_group = detect_duplicates(article_items, DUPLICATE_THRESHOLD)
    else:
        doc_dup_of, doc_dup_score, doc_dup_group = {}, {}, {}
        art_dup_of, art_dup_score, art_dup_group = {}, {}, {}

    iterable2 = tqdm(files, desc="Summarizing") if tqdm else files
    for idx, path in enumerate(iterable2, start=1):
        if progress_cb:
            progress_cb("Summarizing", idx - 1, len(files))
        key = file_key(path)
        doc_id = resume_files.get(key, {}).get("doc_id", f"{run_id}-DOC-{idx:05d}")
        text = raw_texts.get(doc_id, "")
        cached = resume_files.get(key, {})
        extraction_status = extraction_statuses.get(doc_id, "no_text")
        low_text = extraction_status != "ok" or len(text) < MIN_EXTRACTED_CHARS
        if cached.get("doc_summary") and not dry_run:
            summary = cached.get("doc_summary", {})
        elif dry_run or not text.strip() or low_text:
            summary = {}
        else:
            try:
                summary = summarize_document(cfg, text, categories)
                resume_files[key]["doc_summary"] = summary
            except Exception as exc:
                logging.exception("Summarization failed for %s: %s", path, exc)
                errors.append({"stage": "summarize", "file_name": path.name, "file_path": str(path), "error": str(exc)})
                summary = {}

        category = (summary.get("category") or "uncategorized").strip()
        if low_text:
            category = UNREADABLE_CATEGORY
        if category not in categories and category not in ("Other", UNREADABLE_CATEGORY):
            category = "Other"

        tags = summary.get("tags") or []
        if isinstance(tags, str):
            tags = [t.strip() for t in tags.split(",") if t.strip()]
        # normalize/dedupe tags
        seen = set()
        norm_tags = []
        for t in tags:
            t = str(t).strip()
            tag_key = t.lower()
            if not t or tag_key in seen:
                continue
            seen.add(tag_key)
            norm_tags.append(t)
        tags = norm_tags[:MAX_TAGS]

        short_summary = (summary.get("short_summary") or "").strip()
        long_summary = (summary.get("long_summary") or "").strip()

        duplicate_of = doc_dup_of.get(doc_id, "")
        duplicate_score = doc_dup_score.get(doc_id)
        duplicate_group_id = doc_dup_group.get(doc_id, "")

        moved_to = ""
        review_flags = []
        if extraction_status != "ok":
            review_flags.append("low_text")
        if len(text) < MIN_EXTRACTED_CHARS:
            review_flags.append("short_text")

        docs.append(
            DocRecord(
                doc_id=doc_id,
                file_key=key,
                file_name=path.name,
                file_path=str(path),
                file_ext=path.suffix.lower(),
                category=category,
                tags=tags,
                short_summary=short_summary,
                long_summary=long_summary,
                word_count=len(text.split()),
                char_count=len(text),
                extraction_status=extraction_statuses.get(doc_id, "no_text"),
                review_flags=",".join(review_flags),
                duplicate_of=duplicate_of,
                duplicate_score=duplicate_score,
                duplicate_group_id=duplicate_group_id,
                moved_to=moved_to,
            )
        )
        if progress_cb:
            progress_cb("Summarizing", idx, len(files))

        # Article summaries
        # Only for PDF (others have no articles)
        article_list = []
        # We re-split to align with doc order
        if path.suffix.lower() == ".pdf":
            if cached and cached.get("articles_raw"):
                article_list = cached.get("articles_raw", [])
            else:
                try:
                    _, pages, _ = extract_text_pdf(path, ocrmypdf_enabled)
                    article_list = split_pdf_into_articles(pages)
                except Exception as exc:
                    logging.exception("Failed to split articles for %s: %s", path, exc)
                    errors.append({"stage": "split_articles", "file_name": path.name, "file_path": str(path), "error": str(exc)})
        for a_idx, (title, body) in enumerate(article_list, start=1):
            article_id = f"{doc_id}-A{a_idx:03d}"
            cached_summary = None
            if cached and "article_summaries" in cached:
                cached_summary = cached["article_summaries"].get(str(a_idx))
            if cached_summary is not None and not dry_run:
                art_summary = cached_summary
            elif dry_run or not body.strip():
                art_summary = ""
            else:
                try:
                    art_summary = summarize_article(cfg, body)
                    resume_files[key].setdefault("article_summaries", {})[str(a_idx)] = art_summary
                except Exception as exc:
                    logging.exception("Article summary failed for %s: %s", path, exc)
                    errors.append({"stage": "article_summarize", "file_name": path.name, "file_path": str(path), "error": str(exc)})
                    art_summary = ""
            articles.append(
                ArticleRecord(
                    doc_id=doc_id,
                    file_key=key,
                    file_name=path.name,
                    file_path=str(path),
                    article_index=a_idx,
                    article_title=title,
                    article_summary=art_summary,
                    duplicate_of=art_dup_of.get(article_id, ""),
                    duplicate_score=art_dup_score.get(article_id),
                    duplicate_group_id=art_dup_group.get(article_id, ""),
                )
            )

    if embeddings_source == "summary":
        doc_items2: List[Tuple[str, str, Optional[np.ndarray]]] = []
        article_items2: List[Tuple[str, str, Optional[np.ndarray]]] = []
        min_chars = min_embedding_chars_for_source(embeddings_source)

        for d in docs:
            key = doc_id_to_key.get(d.doc_id, "")
            normalized = normalize_text(raw_texts.get(d.doc_id, ""))
            emb_text = embedding_text_for_doc(embeddings_source, normalized, d.long_summary, d.short_summary)
            emb_vec: Optional[np.ndarray] = None
            cached_emb = None
            if key and key in resume_files:
                cached_emb = resume_files[key].get("doc_embedding_summary")
            if cached_emb is not None:
                emb_vec = np.array(cached_emb, dtype=np.float32)
            elif emb_text and len(emb_text) >= min_chars and not dry_run:
                try:
                    emb = call_azure_embeddings(cfg, emb_text[:MAX_CHARS_PER_CHUNK])
                    emb_vec = np.array(emb, dtype=np.float32)
                    if key:
                        resume_files[key]["doc_embedding_summary"] = emb
                except Exception as exc:
                    logging.exception("Embedding failed for %s: %s", d.file_path, exc)
            doc_items2.append((d.doc_id, doc_hashes.get(d.doc_id, ""), emb_vec))

        for a in articles:
            article_id = f"{a.doc_id}-A{a.article_index:03d}"
            key = article_id_to_key.get(article_id, "")
            idx = article_id_to_idx.get(article_id, None)
            emb_text = embedding_text_for_article(embeddings_source, article_texts.get(article_id, ""), a.article_summary)
            aemb_vec: Optional[np.ndarray] = None
            cached_aemb = None
            if key and idx is not None and key in resume_files:
                cached_aemb = resume_files[key].get("article_embedding_summary", {}).get(str(idx))
            if cached_aemb is not None:
                aemb_vec = np.array(cached_aemb, dtype=np.float32)
            elif emb_text and len(emb_text) >= min_chars and not dry_run:
                try:
                    aemb = call_azure_embeddings(cfg, emb_text[:MAX_CHARS_PER_CHUNK])
                    aemb_vec = np.array(aemb, dtype=np.float32)
                    if key and idx is not None:
                        resume_files[key].setdefault("article_embedding_summary", {})[str(idx)] = aemb
                except Exception as exc:
                    logging.exception("Article embedding failed for %s: %s", a.file_path, exc)
            article_items2.append((article_id, article_hashes.get(article_id, ""), aemb_vec))

        doc_dup_of, doc_dup_score, doc_dup_group = detect_duplicates(doc_items2, DUPLICATE_THRESHOLD)
        art_dup_of, art_dup_score, art_dup_group = detect_duplicates(article_items2, DUPLICATE_THRESHOLD)

        for d in docs:
            d.duplicate_of = doc_dup_of.get(d.doc_id, "")
            d.duplicate_score = doc_dup_score.get(d.doc_id)
            d.duplicate_group_id = doc_dup_group.get(d.doc_id, "")

        for a in articles:
            article_id = f"{a.doc_id}-A{a.article_index:03d}"
            a.duplicate_of = art_dup_of.get(article_id, "")
            a.duplicate_score = art_dup_score.get(article_id)
            a.duplicate_group_id = art_dup_group.get(article_id, "")

    # Move files
    if not dry_run and not no_move:
        for i, d in enumerate(docs, start=1):
            if progress_cb:
                progress_cb("Moving files", i - 1, len(docs))
            src = Path(d.file_path)
            cat_folder = sanitize_folder(d.category)
            if d.duplicate_of:
                dest_dir = output_dir / f"{cat_folder}_Duplicate"
            else:
                dest_dir = output_dir / cat_folder
            dest_dir.mkdir(parents=True, exist_ok=True)
            target = dest_dir / src.name
            if target.exists():
                stem = target.stem
                suffix = target.suffix
                i = 1
                while True:
                    candidate = dest_dir / f"{stem}_{i}{suffix}"
                    if not candidate.exists():
                        target = candidate
                        break
                    i += 1
            try:
                shutil.move(str(src), str(target))
                d.moved_to = str(target)
            except Exception as exc:
                logging.exception("Failed to move %s: %s", src, exc)
                errors.append({"stage": "move", "file_name": src.name, "file_path": str(src), "error": str(exc)})
            if progress_cb:
                progress_cb("Moving files", i, len(docs))

    full_text_rows: List[Dict[str, Any]] = []
    for d in docs:
        text = raw_texts.get(d.doc_id, "")
        parts = split_for_excel(text)
        full_text_rows.append(
            {
                "doc_id": d.doc_id,
                "file_key": d.file_key,
                "category": d.category,
                "short_summary": d.short_summary,
                "long_summary": d.long_summary,
                "tags": ", ".join(d.tags),
                "word_count": d.word_count,
                "char_count": d.char_count,
                "extraction_status": d.extraction_status,
                "review_flags": d.review_flags,
                "full_text": parts[0] if parts else "",
                "full_text_parts_count": len(parts),
                "full_text_parts": parts,
                "moved_to": d.moved_to,
                "file_name": d.file_name,
                "file_path": d.file_path,
            }
        )

    # Write outputs
    peers_path = full_text_path = None
    try:
        peers_path, full_text_path = write_excels(output_dir, docs, articles, full_text_rows, app_name, append_excel)
        logging.info("Wrote %s and %s", peers_path, full_text_path)
    except Exception as exc:
        logging.exception("Failed to write Excel outputs: %s", exc)
        errors.append({"stage": "write_excel", "file_name": "", "file_path": str(output_dir), "error": str(exc)})
    usage = get_usage()
    try:
        report_path = write_summary_report(output_dir, docs, articles, errors, usage, total_files, len(files), limit)
        logging.info("Wrote %s", report_path)
    except Exception as exc:
        logging.exception("Failed to write summary report: %s", exc)

    if use_resume:
        resume["files"] = resume_files
        save_resume(output_dir, resume)

    elapsed = time.time() - t0
    save_last_run_stats(
        output_dir,
        {
            "timestamp": time.strftime("%Y-%m-%d %H:%M:%S"),
            "elapsed_sec": elapsed,
            "processed_files": len(files),
            "total_files": total_files,
            "total_size_mb": processed_stats.get("total_size_mb", 0.0),
            "ocr_enabled": ocrmypdf_enabled,
            "embeddings_source": embeddings_source,
            "chat_deployment": cfg.chat_deployment,
        },
    )
    if limit is not None and total_files > len(files):
        est_total = (elapsed / max(len(files), 1)) * total_files
        logging.info(
            "Estimate for %s files based on %s processed: ~%ss",
            total_files,
            len(files),
            int(est_total),
        )
        est_10k = (elapsed / max(len(files), 1)) * 10000
        logging.info("Estimate for 10000 files: ~%ss", int(est_10k))
        est_10k = (elapsed / max(len(files), 1)) * 10000
        logging.info("Estimate for 10000 files: ~%ss", int(est_10k))


def run_pipeline_parallel(
    input_dir: Path,
    output_dir: Path,
    categories: List[str],
    cfg: AzureConfig,
    dry_run: bool,
    use_resume: bool,
    ocrmypdf_enabled: bool,
    app_name: Optional[str],
    embeddings_source: str,
    append_excel: bool,
    workers: int,
    limit: Optional[int] = None,
    no_move: bool = False,
) -> None:
    setup_logging(output_dir)
    logging.info("Starting pipeline (parallel, workers=%s)", workers)

    reset_usage()

    run_id = time.strftime("%Y%m%d%H%M%S")

    files = list_files(input_dir)
    total_files = len(files)
    if limit is not None and limit > 0:
        files = files[:limit]
    processed_stats = scan_input_stats(files)
    if not files:
        logging.warning("No supported files found")
        return

    t0 = time.time()

    resume = load_resume(output_dir) if use_resume else {"files": {}}
    resume_files = resume.get("files", {})

    raw_texts: Dict[str, str] = {}
    article_texts: Dict[str, str] = {}
    extraction_statuses: Dict[str, str] = {}
    doc_items: List[Tuple[str, str, Optional[np.ndarray]]] = []
    article_items: List[Tuple[str, str, Optional[np.ndarray]]] = []
    articles_by_doc: Dict[str, List[Tuple[str, str]]] = {}
    doc_hashes: Dict[str, str] = {}
    article_hashes: Dict[str, str] = {}
    doc_id_to_key: Dict[str, str] = {}
    article_id_to_key: Dict[str, str] = {}
    article_id_to_idx: Dict[str, int] = {}
    errors: List[Dict[str, str]] = []
    errors_lock = threading.Lock()

    def extract_and_embed(idx_path: Tuple[int, Path]) -> Tuple[int, Path, str, List[Tuple[str, str]], str, Optional[np.ndarray], List[Tuple[str, Optional[np.ndarray]]]]:
        idx, path = idx_path
        key = file_key(path)
        cached = resume_files.get(key)
        doc_id = f"{run_id}-DOC-{idx:05d}"
        if cached and cached.get("doc_id"):
            doc_id = cached["doc_id"]
        doc_id_to_key[doc_id] = key
        if cached:
            text = cached.get("text", "")
            pdf_articles = cached.get("articles_raw", [])
            extraction_status = cached.get("extraction_status", "no_text")
            cached["doc_id"] = doc_id
        else:
            try:
                text, pdf_articles, extraction_status = process_file(path, ocrmypdf_enabled)
            except Exception as exc:
                logging.exception("Failed to extract %s: %s", path, exc)
                with errors_lock:
                    errors.append({"stage": "extract", "file_name": path.name, "file_path": str(path), "error": str(exc)})
                text = ""
                pdf_articles = []
                extraction_status = "no_text"
            resume_files[key] = {
                "doc_id": doc_id,
                "file_path": str(path),
                "file_name": path.name,
                "ext": path.suffix.lower(),
                "text": text,
                "articles_raw": pdf_articles,
                "extraction_status": extraction_status,
            }

        normalized = normalize_text(text)
        hsh = hash_text(normalized)

        emb_vec: Optional[np.ndarray] = None
        if embeddings_source == "full_text":
            cached_emb = cached.get("doc_embedding") if cached else None
            if cached_emb is not None:
                emb_vec = np.array(cached_emb, dtype=np.float32)
            elif normalized and len(normalized) >= MIN_EMBEDDING_CHARS and not dry_run:
                try:
                    emb = call_azure_embeddings(cfg, normalized[:MAX_CHARS_PER_CHUNK])
                    emb_vec = np.array(emb, dtype=np.float32)
                    resume_files[key]["doc_embedding"] = emb
                except Exception as exc:
                    logging.exception("Embedding failed for %s: %s", path, exc)

        art_embs: List[Tuple[str, Optional[np.ndarray]]] = []
        for a_idx, (_title, body) in enumerate(pdf_articles, start=1):
            article_id = f"{doc_id}-A{a_idx:03d}"
            article_texts[article_id] = body
            article_id_to_key[article_id] = key
            article_id_to_idx[article_id] = a_idx
            aemb_vec: Optional[np.ndarray] = None
            if embeddings_source == "full_text":
                cached_aemb = None
                if cached and "article_embeddings" in cached:
                    cached_aemb = cached["article_embeddings"].get(str(a_idx))
                if cached_aemb is not None:
                    aemb_vec = np.array(cached_aemb, dtype=np.float32)
                elif body.strip() and len(body) >= MIN_EMBEDDING_CHARS and not dry_run:
                    try:
                        aemb = call_azure_embeddings(cfg, body[:MAX_CHARS_PER_CHUNK])
                        aemb_vec = np.array(aemb, dtype=np.float32)
                        resume_files[key].setdefault("article_embeddings", {})[str(a_idx)] = aemb
                    except Exception as exc:
                        logging.exception("Article embedding failed for %s: %s", path, exc)
                        with errors_lock:
                            errors.append({"stage": "article_embedding", "file_name": path.name, "file_path": str(path), "error": str(exc)})
            art_embs.append((article_id, aemb_vec))

        if use_resume:
            resume["files"] = resume_files
            save_resume(output_dir, resume)

        return idx, path, text, pdf_articles, extraction_status, (hsh, emb_vec), art_embs

    with ThreadPoolExecutor(max_workers=workers) as ex:
        futures = [ex.submit(extract_and_embed, (i, p)) for i, p in enumerate(files, start=1)]
        for fut in as_completed(futures):
            try:
                idx, path, text, pdf_articles, extraction_status, doc_hash_emb, art_embs = fut.result()
            except Exception as exc:
                logging.exception("Worker failed: %s", exc)
                with errors_lock:
                    errors.append({"stage": "worker", "file_name": "", "file_path": "", "error": str(exc)})
                continue
            doc_id = resume_files.get(file_key(path), {}).get("doc_id", f"{run_id}-DOC-{idx:05d}")
            raw_texts[doc_id] = text
            extraction_statuses[doc_id] = extraction_status
            hsh, emb_vec = doc_hash_emb
            doc_hashes[doc_id] = hsh
            doc_items.append((doc_id, hsh, emb_vec))
            articles_by_doc[doc_id] = pdf_articles
            for (article_id, aemb_vec), (title, body) in zip(art_embs, pdf_articles):
                ahash = hash_text(normalize_text(body))
                article_hashes[article_id] = ahash
                article_items.append((article_id, ahash, aemb_vec))

    if embeddings_source == "full_text":
        doc_dup_of, doc_dup_score, doc_dup_group = detect_duplicates(doc_items, DUPLICATE_THRESHOLD)
        art_dup_of, art_dup_score, art_dup_group = detect_duplicates(article_items, DUPLICATE_THRESHOLD)
    elif embeddings_source == EMBEDDINGS_SOURCE_NONE:
        doc_dup_of, doc_dup_score, doc_dup_group = detect_duplicates(doc_items, DUPLICATE_THRESHOLD)
        art_dup_of, art_dup_score, art_dup_group = detect_duplicates(article_items, DUPLICATE_THRESHOLD)
    else:
        doc_dup_of, doc_dup_score, doc_dup_group = {}, {}, {}
        art_dup_of, art_dup_score, art_dup_group = {}, {}, {}

    docs: List[DocRecord] = []
    articles: List[ArticleRecord] = []

    def summarize_doc(idx_path: Tuple[int, Path]) -> Tuple[int, Path, Dict[str, Any]]:
        idx, path = idx_path
        key = file_key(path)
        doc_id = resume_files.get(key, {}).get("doc_id", f"{run_id}-DOC-{idx:05d}")
        text = raw_texts.get(doc_id, "")
        cached = resume_files.get(key, {})
        extraction_status = extraction_statuses.get(doc_id, "no_text")
        low_text = extraction_status != "ok" or len(text) < MIN_EXTRACTED_CHARS
        if cached.get("doc_summary") and not dry_run:
            summary = cached.get("doc_summary", {})
        elif dry_run or not text.strip() or low_text:
            summary = {}
        else:
            try:
                summary = summarize_document(cfg, text, categories)
                resume_files[key]["doc_summary"] = summary
            except Exception as exc:
                logging.exception("Summarization failed for %s: %s", path, exc)
                with errors_lock:
                    errors.append({"stage": "summarize", "file_name": path.name, "file_path": str(path), "error": str(exc)})
                summary = {}
        if use_resume:
            resume["files"] = resume_files
            save_resume(output_dir, resume)
        return idx, path, summary

    with ThreadPoolExecutor(max_workers=workers) as ex:
        futures = [ex.submit(summarize_doc, (i, p)) for i, p in enumerate(files, start=1)]
        summaries: Dict[str, Dict[str, Any]] = {}
        for fut in as_completed(futures):
            try:
                idx, path, summary = fut.result()
            except Exception as exc:
                logging.exception("Summarize worker failed: %s", exc)
                with errors_lock:
                    errors.append({"stage": "summarize_worker", "file_name": "", "file_path": "", "error": str(exc)})
                continue
            doc_id = f"DOC-{idx:05d}"
            summaries[doc_id] = summary

    for idx, path in enumerate(files, start=1):
        key = file_key(path)
        doc_id = resume_files.get(key, {}).get("doc_id", f"{run_id}-DOC-{idx:05d}")
        summary = summaries.get(doc_id, {})
        text = raw_texts.get(doc_id, "")
        extraction_status = extraction_statuses.get(doc_id, "no_text")
        low_text = extraction_status != "ok" or len(text) < MIN_EXTRACTED_CHARS

        category = (summary.get("category") or "uncategorized").strip()
        if low_text:
            category = UNREADABLE_CATEGORY
        if category not in categories and category not in ("Other", UNREADABLE_CATEGORY):
            category = "Other"

        tags = summary.get("tags") or []
        if isinstance(tags, str):
            tags = [t.strip() for t in tags.split(",") if t.strip()]
        seen = set()
        norm_tags = []
        for t in tags:
            t = str(t).strip()
            tag_key = t.lower()
            if not t or tag_key in seen:
                continue
            seen.add(tag_key)
            norm_tags.append(t)
        tags = norm_tags[:MAX_TAGS]

        short_summary = (summary.get("short_summary") or "").strip()
        long_summary = (summary.get("long_summary") or "").strip()

        duplicate_of = doc_dup_of.get(doc_id, "")
        duplicate_score = doc_dup_score.get(doc_id)
        duplicate_group_id = doc_dup_group.get(doc_id, "")

        review_flags = []
        if extraction_status != "ok":
            review_flags.append("low_text")
        if len(text) < MIN_EXTRACTED_CHARS:
            review_flags.append("short_text")

        docs.append(
            DocRecord(
                doc_id=doc_id,
                file_key=key,
                file_name=path.name,
                file_path=str(path),
                file_ext=path.suffix.lower(),
                category=category,
                tags=tags,
                short_summary=short_summary,
                long_summary=long_summary,
                word_count=len(text.split()),
                char_count=len(text),
                extraction_status=extraction_statuses.get(doc_id, "no_text"),
                review_flags=",".join(review_flags),
                duplicate_of=duplicate_of,
                duplicate_score=duplicate_score,
                duplicate_group_id=duplicate_group_id,
                moved_to="",
            )
        )

        article_list = articles_by_doc.get(doc_id, [])
        for a_idx, (title, body) in enumerate(article_list, start=1):
            article_id = f"{doc_id}-A{a_idx:03d}"
            cached = resume_files.get(key, {})
            cached_summary = None
            if cached and "article_summaries" in cached:
                cached_summary = cached["article_summaries"].get(str(a_idx))
            if cached_summary is not None and not dry_run:
                art_summary = cached_summary
            elif dry_run or not body.strip():
                art_summary = ""
            else:
                try:
                    art_summary = summarize_article(cfg, body)
                    resume_files[key].setdefault("article_summaries", {})[str(a_idx)] = art_summary
                except Exception as exc:
                    logging.exception("Article summary failed for %s: %s", path, exc)
                    with errors_lock:
                        errors.append({"stage": "article_summarize", "file_name": path.name, "file_path": str(path), "error": str(exc)})
                    art_summary = ""
            articles.append(
                ArticleRecord(
                    doc_id=doc_id,
                    file_key=key,
                    file_name=path.name,
                    file_path=str(path),
                    article_index=a_idx,
                    article_title=title,
                    article_summary=art_summary,
                    duplicate_of=art_dup_of.get(article_id, ""),
                    duplicate_score=art_dup_score.get(article_id),
                    duplicate_group_id=art_dup_group.get(article_id, ""),
                )
            )

    if embeddings_source == "summary":
        doc_items2: List[Tuple[str, str, Optional[np.ndarray]]] = []
        article_items2: List[Tuple[str, str, Optional[np.ndarray]]] = []
        min_chars = min_embedding_chars_for_source(embeddings_source)

        for d in docs:
            key = doc_id_to_key.get(d.doc_id, "")
            normalized = normalize_text(raw_texts.get(d.doc_id, ""))
            emb_text = embedding_text_for_doc(embeddings_source, normalized, d.long_summary, d.short_summary)
            emb_vec: Optional[np.ndarray] = None
            cached_emb = None
            if key and key in resume_files:
                cached_emb = resume_files[key].get("doc_embedding_summary")
            if cached_emb is not None:
                emb_vec = np.array(cached_emb, dtype=np.float32)
            elif emb_text and len(emb_text) >= min_chars and not dry_run:
                try:
                    emb = call_azure_embeddings(cfg, emb_text[:MAX_CHARS_PER_CHUNK])
                    emb_vec = np.array(emb, dtype=np.float32)
                    if key:
                        resume_files[key]["doc_embedding_summary"] = emb
                except Exception as exc:
                    logging.exception("Embedding failed for %s: %s", d.file_path, exc)
            doc_items2.append((d.doc_id, doc_hashes.get(d.doc_id, ""), emb_vec))

        for a in articles:
            article_id = f"{a.doc_id}-A{a.article_index:03d}"
            key = article_id_to_key.get(article_id, "")
            idx = article_id_to_idx.get(article_id, None)
            emb_text = embedding_text_for_article(embeddings_source, article_texts.get(article_id, ""), a.article_summary)
            aemb_vec: Optional[np.ndarray] = None
            cached_aemb = None
            if key and idx is not None and key in resume_files:
                cached_aemb = resume_files[key].get("article_embedding_summary", {}).get(str(idx))
            if cached_aemb is not None:
                aemb_vec = np.array(cached_aemb, dtype=np.float32)
            elif emb_text and len(emb_text) >= min_chars and not dry_run:
                try:
                    aemb = call_azure_embeddings(cfg, emb_text[:MAX_CHARS_PER_CHUNK])
                    aemb_vec = np.array(aemb, dtype=np.float32)
                    if key and idx is not None:
                        resume_files[key].setdefault("article_embedding_summary", {})[str(idx)] = aemb
                except Exception as exc:
                    logging.exception("Article embedding failed for %s: %s", a.file_path, exc)
            article_items2.append((article_id, article_hashes.get(article_id, ""), aemb_vec))

        doc_dup_of, doc_dup_score, doc_dup_group = detect_duplicates(doc_items2, DUPLICATE_THRESHOLD)
        art_dup_of, art_dup_score, art_dup_group = detect_duplicates(article_items2, DUPLICATE_THRESHOLD)

        for d in docs:
            d.duplicate_of = doc_dup_of.get(d.doc_id, "")
            d.duplicate_score = doc_dup_score.get(d.doc_id)
            d.duplicate_group_id = doc_dup_group.get(d.doc_id, "")

        for a in articles:
            article_id = f"{a.doc_id}-A{a.article_index:03d}"
            a.duplicate_of = art_dup_of.get(article_id, "")
            a.duplicate_score = art_dup_score.get(article_id)
            a.duplicate_group_id = art_dup_group.get(article_id, "")

    if not dry_run and not no_move:
        for d in docs:
            src = Path(d.file_path)
            cat_folder = sanitize_folder(d.category)
            if d.duplicate_of:
                dest_dir = output_dir / f"{cat_folder}_Duplicate"
            else:
                dest_dir = output_dir / cat_folder
            dest_dir.mkdir(parents=True, exist_ok=True)
            target = dest_dir / src.name
            if target.exists():
                stem = target.stem
                suffix = target.suffix
                i = 1
                while True:
                    candidate = dest_dir / f"{stem}_{i}{suffix}"
                    if not candidate.exists():
                        target = candidate
                        break
                    i += 1
            try:
                shutil.move(str(src), str(target))
                d.moved_to = str(target)
            except Exception as exc:
                logging.exception("Failed to move %s: %s", src, exc)
                with errors_lock:
                    errors.append({"stage": "move", "file_name": src.name, "file_path": str(src), "error": str(exc)})

    full_text_rows: List[Dict[str, Any]] = []
    for d in docs:
        text = raw_texts.get(d.doc_id, "")
        parts = split_for_excel(text)
        full_text_rows.append(
            {
                "doc_id": d.doc_id,
                "file_key": d.file_key,
                "category": d.category,
                "short_summary": d.short_summary,
                "long_summary": d.long_summary,
                "tags": ", ".join(d.tags),
                "word_count": d.word_count,
                "char_count": d.char_count,
                "extraction_status": d.extraction_status,
                "full_text": parts[0] if parts else "",
                "full_text_parts_count": len(parts),
                "full_text_parts": parts,
                "moved_to": d.moved_to,
                "file_name": d.file_name,
                "file_path": d.file_path,
            }
        )

    peers_path = full_text_path = None
    try:
        peers_path, full_text_path = write_excels(output_dir, docs, articles, full_text_rows, app_name, append_excel)
        logging.info("Wrote %s and %s", peers_path, full_text_path)
    except Exception as exc:
        logging.exception("Failed to write Excel outputs: %s", exc)
        with errors_lock:
            errors.append({"stage": "write_excel", "file_name": "", "file_path": str(output_dir), "error": str(exc)})
    usage = get_usage()
    try:
        report_path = write_summary_report(output_dir, docs, articles, errors, usage, total_files, len(files), limit)
        logging.info("Wrote %s", report_path)
    except Exception as exc:
        logging.exception("Failed to write summary report: %s", exc)

    if use_resume:
        resume["files"] = resume_files
        save_resume(output_dir, resume)

    elapsed = time.time() - t0
    save_last_run_stats(
        output_dir,
        {
            "timestamp": time.strftime("%Y-%m-%d %H:%M:%S"),
            "elapsed_sec": elapsed,
            "processed_files": len(files),
            "total_files": total_files,
            "total_size_mb": processed_stats.get("total_size_mb", 0.0),
            "ocr_enabled": ocrmypdf_enabled,
            "embeddings_source": embeddings_source,
            "chat_deployment": cfg.chat_deployment,
        },
    )
    if limit is not None and total_files > len(files):
        est_total = (elapsed / max(len(files), 1)) * total_files
        logging.info(
            "Estimate for %s files based on %s processed: ~%ss",
            total_files,
            len(files),
            int(est_total),
        )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="DocAtlas document processing pipeline")
    parser.add_argument("--input", help="Input folder")
    parser.add_argument("--output", help="Output folder")
    parser.add_argument("--categories", help="Categories separated by semicolons")
    parser.add_argument("--config", help="Path to applications config JSON")
    parser.add_argument("--app", help="Application name from config")
    parser.add_argument("--dry-run", action="store_true", help="Do not call APIs or move files")
    parser.add_argument("--no-resume", action="store_true", help="Disable resume cache")
    parser.add_argument("--no-ocrmypdf", action="store_true", help="Disable OCRmyPDF and use Tesseract fallback")
    parser.add_argument("--edit-config", action="store_true", help="Open GUI editor for applications config")
    parser.add_argument("--embeddings-source", choices=["summary", "full_text", "none"], help="Use summaries, full text, or disable embeddings")
    parser.add_argument("--overwrite-excel", action="store_true", help="Overwrite Excel outputs instead of appending")
    parser.add_argument("--limit", type=int, help="Process only the first N files (for estimation)")
    parser.add_argument("--no-move", action="store_true", help="Do not move files (for estimation)")
    parser.add_argument("--charter-mode", action="store_true", help="Preview-only mode (no file moves)")
    parser.add_argument("--signal-scan", action="store_true", help="Deprecated alias for --charter-mode")
    parser.add_argument("--test-embeddings", action="store_true", help="Test embeddings endpoint and exit")
    parser.add_argument("--test-chat", action="store_true", help="Test chat endpoint and exit")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    config_path = Path(args.config) if args.config else Path(__file__).with_name("applications.json")
    app_config = load_app_config(config_path)
    is_gui_flow = not (args.input and args.output and (args.categories or args.app))

    if args.edit_config:
        if tk is None:
            raise RuntimeError("tkinter is not available")
        root = tk.Tk()
        root.withdraw()
        edit_applications_gui(config_path, app_config, root)
        root.mainloop()
        return 0

    if args.test_embeddings:
        cfg = azure_config_from_env(require_key=False)
        if not cfg.embeddings_api_key:
            cfg.embeddings_api_key = prompt_api_key_gui(
                "DocAtlas - Enter Embeddings API Key",
                "Paste embeddings API key (not stored):",
            ) or ""
        if not cfg.embeddings_api_key:
            raise ValueError("AZURE_EMBEDDINGS_API_KEY is not set")
        try:
            emb = call_azure_embeddings(cfg, "test embedding")
            print(f"Embeddings OK. Vector length: {len(emb)}")
            return 0
        except Exception as exc:
            print(f"Embeddings test failed: {exc}")
            return 1

    if args.test_chat:
        cfg = azure_config_from_env(require_key=False)
        if not cfg.chat_api_key:
            cfg.chat_api_key = prompt_api_key_gui(
                "DocAtlas - Enter LLM API Key",
                "Paste LLM API key (not stored):",
            ) or ""
        if not cfg.chat_api_key:
            raise ValueError("AZURE_CHAT_API_KEY is not set")
        if not cfg.api_key:
            cfg.api_key = cfg.chat_api_key
        try:
            msg = [{"role": "user", "content": "Say OK"}]
            out = call_azure_chat(cfg, msg)
            print("Chat OK. Response:", out[:200])
            return 0
        except Exception as exc:
            print(f"Chat test failed: {exc}")
            return 1

    append_excel = not args.overwrite_excel
    if args.charter_mode or args.signal_scan:
        args.no_move = True

    if args.input and args.output and (args.categories or args.app):
        input_dir = Path(args.input)
        output_dir = Path(args.output)
        if args.categories:
            categories = [c.strip() for c in args.categories.split(";") if c.strip()]
            app_name = None
        elif args.app and args.app in app_config:
            categories = app_config[args.app]
            app_name = args.app
        else:
            raise ValueError("Provide --categories or a valid --app from config")
        ocrmypdf_enabled = not args.no_ocrmypdf
        embeddings_source = resolve_embeddings_source(args.embeddings_source)
    else:
        input_dir, output_dir = pick_directories_gui()
        categories, app_name = get_categories_gui(app_config, config_path)
        ocrmypdf_enabled = get_ocrmypdf_gui()
        embeddings_source, append_excel = get_embeddings_source_gui()
        gui_charter_mode = get_run_mode_gui()
        if gui_charter_mode is None:
            return 0
        if gui_charter_mode:
            args.no_move = True

    cfg = azure_config_from_env(require_key=(not args.dry_run and not is_gui_flow))
    if not args.dry_run:
        if not cfg.chat_api_key:
            if is_gui_flow:
                cfg.chat_api_key = prompt_api_key_gui("DocAtlas - Enter LLM API Key", "Paste LLM API key (not stored):") or ""
            if not cfg.chat_api_key:
                raise ValueError("AZURE_CHAT_API_KEY is not set")
        if embeddings_source != EMBEDDINGS_SOURCE_NONE and not cfg.embeddings_api_key:
            if is_gui_flow:
                cfg.embeddings_api_key = prompt_api_key_gui(
                    "DocAtlas - Enter Embeddings API Key",
                    "Paste embeddings API key (not stored):",
                ) or ""
            if embeddings_source != EMBEDDINGS_SOURCE_NONE and not cfg.embeddings_api_key:
                raise ValueError("AZURE_EMBEDDINGS_API_KEY is not set")
        if not cfg.api_key:
            cfg.api_key = cfg.chat_api_key
    warn_missing_ocr_deps(ocrmypdf_enabled)
    try:
        est_files = list_files(input_dir)
        est_stats = scan_input_stats(est_files)
        est_sec, est_source, settings_match = quick_estimate_runtime(
            est_stats,
            output_dir,
            ocrmypdf_enabled,
            embeddings_source,
            cfg.chat_deployment,
        )
        if is_gui_flow and messagebox is not None and est_sec:
            note = ""
            if est_source == "baseline" and not settings_match:
                note = "\nNote: baseline settings differ from current run."
            msg = (
                f"Files: {est_stats.get('count', 0)}\n"
                f"Total size: {est_stats.get('total_size_mb', 0.0):.1f} MB\n"
                f"Estimated time: ~{format_duration(est_sec)}{note}"
            )
            messagebox.showinfo("DocAtlas - Estimated Runtime", msg)
    except Exception:
        pass

    if not (args.input and args.output and (args.categories or args.app)):
        # GUI progress window
        if tk is None:
            raise RuntimeError("tkinter is not available")

        q: "queue.Queue[Tuple[str, int, int]]" = queue.Queue()
        stage_times: Dict[str, float] = {}
        stage_start: Dict[str, float] = {}

        def progress_cb(stage: str, current: int, total: int) -> None:
            q.put((stage, current, total))

        def worker() -> None:
            run_pipeline(
                input_dir,
                output_dir,
                categories,
                cfg,
                args.dry_run,
                not args.no_resume,
                ocrmypdf_enabled,
                app_name,
                embeddings_source,
                append_excel,
                args.limit,
                args.no_move,
                progress_cb,
            )
            q.put(("DONE", 1, 1))

        progress_root = tk.Tk()
        progress_root.title("DocAtlas - Processing Documents")
        progress_root.geometry("520x180")
        apply_theme(progress_root)

        frame = tk.Frame(progress_root, bg=THEME["bg"])
        frame.pack(fill=tk.BOTH, expand=True, padx=16, pady=16)

        stage_label = tk.Label(frame, text="Starting...", bg=THEME["bg"], fg=THEME["fg"], font=FONT_HEADER)
        stage_label.pack(anchor="w", pady=(0, 6))

        eta_label = tk.Label(frame, text="", bg=THEME["bg"], fg=THEME["muted"], font=FONT_SMALL)
        eta_label.pack(anchor="w", pady=(0, 10))

        prog = ttk.Progressbar(frame, length=460, mode="determinate")
        prog.pack(pady=6)

        def poll() -> None:
            try:
                while True:
                    stage, current, total = q.get_nowait()
                    if stage == "DONE":
                        stage_label.config(text="Completed")
                        eta_label.config(text="Finished")
                        prog["value"] = 100
                        progress_root.after(1200, progress_root.destroy)
                        return

                    if stage not in stage_start:
                        stage_start[stage] = time.time()
                    elapsed = time.time() - stage_start[stage]
                    if current > 0:
                        eta = (elapsed / current) * (total - current)
                        eta_label.config(text=f"ETA: ~{int(eta)}s")
                    else:
                        eta_label.config(text="ETA: estimating...")

                    stage_label.config(text=f"{stage} ({current}/{total})")
                    prog["maximum"] = max(total, 1)
                    prog["value"] = current
            except queue.Empty:
                pass
            progress_root.after(200, poll)

        t = threading.Thread(target=worker, daemon=True)
        t.start()
        poll()
        progress_root.mainloop()
    else:
        run_pipeline(
            input_dir,
            output_dir,
            categories,
            cfg,
            args.dry_run,
            not args.no_resume,
            ocrmypdf_enabled,
            app_name,
            embeddings_source,
            append_excel,
            args.limit,
            args.no_move,
        )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
